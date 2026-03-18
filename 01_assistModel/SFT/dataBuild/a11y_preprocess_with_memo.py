#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PPTX/PDF 접근성 진단 보고서 전처리 파이프라인 (OpenAI VLM)

핵심 동작(요구사항 반영):
1) 입력(.pptx 또는 .pdf)에서 "존재하는 모든 이미지"를 추출하고, 각 이미지의 페이지(슬라이드) 내 bbox를 확보
2) VLM으로 각 이미지를 분류: error_region / error_code 
3) error_code 이미지가 있으면 반드시 텍스트로 변환(코드 추출)
4) 페이지(슬라이드) 전체 스크린샷에서 "모든 이미지 영역"을 마스킹 후,
   검사항목/오류유형/평가근거/개선방안코드/오류코드(텍스트일 때만) 를 "보이는대로" 추출
5) 오류코드 텍스트 우선순위: 슬라이드/페이지 텍스트 > 이미지에서 추출한 코드 텍스트(들)

설치:
  pip install python-pptx pillow openai pydantic pymupdf

시스템 의존(PPTX 입력일 때만 필요):
  LibreOffice (soffice)  # PPTX -> PDF 변환에 사용

실행:
  set OPENAI_API_KEY=...    (Windows)
  export OPENAI_API_KEY=... (Linux/macOS)

  python a11y_preprocess.py --input report.pptx --out ./out
  python a11y_preprocess.py --input report.pdf  --out ./out
"""

from __future__ import annotations

import argparse
import base64
import dataclasses
import hashlib
import json
import mimetypes
import os
import platform
import re
import shutil
import subprocess
import sys
import traceback
import zipfile
import posixpath
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple

from PIL import Image, ImageDraw

from pydantic import BaseModel, Field
from openai import OpenAI

import fitz  # PyMuPDF

# python-pptx는 PPTX 입력에서만 사용
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    Presentation = None  # type: ignore
    MSO_SHAPE_TYPE = None  # type: ignore


# -----------------------------
# 설정
# -----------------------------
DEFAULT_MODEL = os.getenv("VLM_MODEL", "gpt-4o")
DEFAULT_TEMPERATURE = float(os.getenv("VLM_TEMPERATURE", "0.0"))
DEFAULT_DPI = int(os.getenv("RENDER_DPI", "200"))

Label = Literal["error_region", "error_code"]


# -----------------------------
# Structured Outputs용 스키마
# -----------------------------
class ImageClassification(BaseModel):
    label: Label


class CodeExtraction(BaseModel):
    code: str = ""


class SlideFields(BaseModel):
    inspection_item: Optional[str] = None
    error_type: Optional[str] = None
    rationale_bullets: List[str] = Field(default_factory=list)
    fix_code: Optional[str] = None
    error_code_text: Optional[str] = None


# -----------------------------
# 데이터 구조
# -----------------------------
@dataclass
class ImageAsset:
    page_index: int               # 슬라이드/페이지 인덱스
    image_index: int              # 페이지 내 이미지 인덱스
    path: Path                    # 추출 이미지 파일 경로
    sha1: str
    ext: str
    bbox_px: Tuple[int, int, int, int]  # x, y, w, h (렌더링 PNG 기준 px)
    label: Optional[Label] = None


@dataclass
class PageRecord:
    page_index: int
    page_image_path: Path
    page_masked_path: Path
    error_region_images: List[Path]      # 최종 산출물에 포함(이미지)
    error_code_images: List[Path]        # 참고용(원본 이미지 경로) - 최종은 텍스트화
    inspection_item: Optional[str]
    error_type: Optional[str]
    rationale_bullets: List[str]
    fix_code: Optional[str]
    error_code_text: Optional[str]
    # provenance(선택)
    error_code_text_from_page: Optional[str]
    error_code_texts_from_images: List[str]



    memo: Optional[str] = None
# -----------------------------
# 유틸
# -----------------------------

# -----------------------------
# Progress bar (optional)
# -----------------------------
_PROGRESS_ENABLED = True
try:
    from tqdm import tqdm as _tqdm  # type: ignore
except Exception:
    _tqdm = None  # type: ignore

def pbar(iterable, **kwargs):
    """
    tqdm이 있으면 진행바를 사용하고, 없거나 비활성화면 원본 iterable을 그대로 반환.
    """
    if (not _PROGRESS_ENABLED) or (_tqdm is None):
        return iterable
    return _tqdm(iterable, **kwargs)

def stage(msg: str) -> None:
    # tqdm 사용/미사용 모두에서 단계 경계를 명확히 보여주기 위한 로그
    print(f"\n== {msg} ==")

def ensure_dir(p: Path) -> None:
    p.mkdir(parents=True, exist_ok=True)
# -----------------------------

def sha1_bytes(b: bytes) -> str:
    return hashlib.sha1(b).hexdigest()


def safe_json_loads(s: str) -> Any:
    s = s.strip()
    s = re.sub(r"^```(json)?\s*", "", s)
    s = re.sub(r"\s*```$", "", s)

    try:
        return json.loads(s)
    except Exception:
        pass

    i, j = s.find("{"), s.rfind("}")
    if i != -1 and j != -1 and j > i:
        return json.loads(s[i : j + 1])

    i, j = s.find("["), s.rfind("]")
    if i != -1 and j != -1 and j > i:
        return json.loads(s[i : j + 1])

    raise ValueError("JSON 파싱 실패: 반환 문자열에 JSON이 없습니다.")


def _guess_mime(path: Path) -> str:
    mime, _ = mimetypes.guess_type(str(path))
    if mime:
        return mime
    ext = path.suffix.lower().lstrip(".")
    if ext in ("jpg", "jpeg"):
        return "image/jpeg"
    return "image/png"


def _data_url_from_image(path: Path) -> str:
    mime = _guess_mime(path)
    b64 = base64.b64encode(path.read_bytes()).decode("utf-8")
    return f"data:{mime};base64,{b64}"


def _run_cmd(cmd: List[str]) -> None:
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(
            "Command failed:\n"
            f"  cmd: {' '.join(cmd)}\n"
            f"  stdout: {proc.stdout}\n"
            f"  stderr: {proc.stderr}\n"
        )


def _find_soffice() -> str:
    """
    LibreOffice 실행 파일(soffice)을 찾는다.
    우선순위:
      1) 환경변수 SOFFICE_PATH (또는 LIBREOFFICE_PATH)
      2) PATH (shutil.which)
      3) OS별 기본 설치 경로 후보
    """
    # 1) env var로 명시
    for k in ("SOFFICE_PATH", "LIBREOFFICE_PATH"):
        v = os.getenv(k, "").strip().strip('"')
        if v and Path(v).exists():
            return v

    # 2) PATH에서 탐색
    for candidate in ("soffice", "soffice.exe"):
        p = shutil.which(candidate)
        if p:
            return p

    # 3) OS별 흔한 설치 경로
    system = platform.system().lower()
    candidates = []

    if system.startswith("windows"):
        candidates += [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            str(Path(os.getenv("LOCALAPPDATA", "")) / "Programs" / "LibreOffice" / "program" / "soffice.exe"),
        ]
    elif system == "darwin":
        candidates += [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]
    else:
        candidates += [
            "/usr/bin/soffice",
            "/usr/local/bin/soffice",
            "/snap/bin/libreoffice",
        ]

    for p in candidates:
        if p and Path(p).exists():
            return p

    raise RuntimeError(
        "LibreOffice(soffice)를 찾지 못했습니다. PPTX 입력을 처리하려면 LibreOffice를 설치하고,\n"
        "1) PATH에 추가하거나, 2) SOFFICE_PATH 환경변수에 soffice 경로를 지정하세요.\n"
        "PDF 입력은 LibreOffice 없이도 동작합니다."
    )



# -----------------------------
# OpenAI VLM 백엔드
# -----------------------------
class OpenAIVLMBackend:
    """
    OpenAI Responses API 기반 VLM 호출.

    - 가능한 경우 responses.parse(Structured Outputs) 사용
    - SDK/버전 차이로 parse가 불가능하면 create + JSON 파싱 fallback
    """

    def __init__(self, model: str = DEFAULT_MODEL, temperature: float = DEFAULT_TEMPERATURE, api_key: Optional[str] = None):
        api_key = 'YOUR API KEY'
        if not api_key:
            raise RuntimeError("OPENAI_API_KEY가 설정되어 있어야 합니다. (환경변수로 설정 권장)")
        self.client = OpenAI(api_key=api_key)
        self.model = model
        self.temperature = temperature
        self.debug = os.getenv("VLM_DEBUG", "").strip().lower() in ("1", "true", "yes", "y", "on")

    def _responses_parse_available(self) -> bool:
        return hasattr(self.client, "responses") and hasattr(self.client.responses, "parse")

    def _log(self, msg: str) -> None:
        if self.debug:
            print(msg)

    def _pydantic_json_schema(self, schema_model) -> Dict[str, Any]:
        """
        pydantic v2: model_json_schema()
        pydantic v1: schema()
        """
        try:
            if hasattr(schema_model, "model_json_schema"):
                return schema_model.model_json_schema()  # type: ignore
            if hasattr(schema_model, "schema"):
                return schema_model.schema()  # type: ignore
        except Exception:
            pass
        return {}

    def _call_parse_or_fallback(self, *, schema_model, input_payload: List[Dict[str, Any]]) -> Any:
        parse_err: Optional[BaseException] = None
        if self._responses_parse_available():
            try:
                resp = self.client.responses.parse(
                    model=self.model,
                    input=input_payload,
                    temperature=self.temperature,
                    text_format=schema_model,
                )
                return resp.output_parsed
            except Exception as e:
                parse_err = e
                self._log(f"[VLM_DEBUG] responses.parse failed: {type(e).__name__}: {e}")
                self._log(traceback.format_exc())

        # 2차 시도: responses.create에서도 JSON schema 강제(지원 시)
        create_kwargs: Dict[str, Any] = {}
        schema = self._pydantic_json_schema(schema_model)
        if schema:
            # OpenAI Responses JSON schema 포맷(지원되는 SDK/버전에서만 동작)
            create_kwargs["text"] = {
                "format": {
                    "type": "json_schema",
                    "name": getattr(schema_model, "__name__", "Schema"),
                    "schema": schema,
                    "strict": True,
                }
            }
        try:
            resp = self.client.responses.create(
                model=self.model,
                input=input_payload,
                temperature=self.temperature,
                **create_kwargs,
            )
        except TypeError as e:
            # 구버전 SDK 등에서 text=... 인자 미지원 시, 일반 create로 폴백
            self._log(f"[VLM_DEBUG] responses.create(text=...) not supported: {e}")
            resp = self.client.responses.create(
                model=self.model,
                input=input_payload,
                temperature=self.temperature,
            )
        except Exception as e:
            # API 단 에러는 여기서 바로 죽이지 말고 “빈 결과”로 처리
            self._log(f"[VLM_DEBUG] responses.create failed: {type(e).__name__}: {e}")
            self._log(traceback.format_exc())
            return {
                "_raw_text": "",
                "_parse_error": f"responses.create failed: {type(e).__name__}: {e}",
                "_parse_stage": "create",
                "_parse_stage_prev": f"parse failed: {type(parse_err).__name__}: {parse_err}" if parse_err else None,
            }


        text = getattr(resp, "output_text", None) or ""
        if not text:
            text = str(resp)

        try:
            return safe_json_loads(text)
        except Exception as e:
            # “JSON이 없다” 케이스 포함: 크래시 금지, raw로 넘겨서 상위 로직이 null 처리하도록
            self._log(f"[VLM_DEBUG] safe_json_loads failed: {type(e).__name__}: {e}")
            self._log(f"[VLM_DEBUG] raw output (truncated 500): {text[:500]!r}")
            return {
                "_raw_text": text,
                "_parse_error": f"{type(e).__name__}: {e}",
                "_parse_stage": "safe_json_loads",
                "_parse_stage_prev": f"parse failed: {type(parse_err).__name__}: {parse_err}" if parse_err else None,
            }

    def classify_image(self, image_path: Path) -> Label:
        prompt = (
            "너는 접근성 진단 보고서(슬라이드/페이지)에서 추출한 이미지를 분류한다.\n"
            "다음 2개 중 하나로만 분류하라.\n"
            "1) error_region: 접근성 오류 상황 UI(검은 화면/페이지 화면/강조 박스 등)\n"
            "2) error_code: HTML/CSS/JS 또는 DOM/Inspector 코드 캡처(코드 스니펫 이미지)\n"
            "출력은 반드시 JSON 객체 1개만: {\"label\":\"error_region\"|\"error_code\"}"
            "예를 들어, {\"label\":\"error_region\"}\n"
            "주의: label 값은 \"error_region\" 또는 \"error_code\" 중 하나."
        )
        input_payload = [{
            "role": "user",
            "content": [
                {"type": "input_text", "text": prompt},
                {"type": "input_image", "image_url": _data_url_from_image(image_path)},
            ],
        }]

        parsed = self._call_parse_or_fallback(schema_model=ImageClassification, input_payload=input_payload)

        if isinstance(parsed, ImageClassification):
            return parsed.label

        if isinstance(parsed, dict):
            label = parsed.get("label")
            if label in ("error_region", "error_code"):
                return label
        return "error_region"

    def extract_code_from_image(self, image_path: Path) -> str:
        prompt = (
            "이미지 안에 보이는 코드(HTML/CSS/JS/DOM 등)를 가능한 한 원문 그대로 텍스트로 추출하라.\n"
            "- 줄바꿈/들여쓰기/기호를 최대한 유지\n"
            "- 설명/해석/요약 금지\n"
            "출력은 반드시 JSON 객체 1개만:\n"
            "{\"code\":\"\"}\n"
            "주의: code는 문자열이며, 줄바꿈/따옴표/역슬래시는 JSON 문자열 규칙에 맞게 이스케이프하라."
        )
        input_payload = [{
            "role": "user",
            "content": [
                {"type": "input_text", "text": prompt},
                {"type": "input_image", "image_url": _data_url_from_image(image_path)},
            ],
        }]

        parsed = self._call_parse_or_fallback(schema_model=CodeExtraction, input_payload=input_payload)

        if isinstance(parsed, CodeExtraction):
            return (parsed.code or "").strip()
        if isinstance(parsed, dict):
            return str(parsed.get("code", "")).strip()
        return ""

    def extract_fields_from_masked_page(self, masked_page_path: Path) -> SlideFields:
        prompt = (
            "너는 접근성 진단 보고서의 '상세 진단 결과' 페이지(슬라이드/페이지)에서 텍스트 필드를 추출한다.\n"
            "주의: 이 이미지는 전체 스크린샷이며, 이미지 영역(오류영역/오류코드/로고 등)은 마스킹되어 있다.\n"
            "따라서 마스킹 영역 내부 내용은 절대 추출하지 말고, 오직 보이는 텍스트만으로 아래 필드를 채워라.\n\n"
            "필드:\n"
            "- inspection_item: 검사항목\n"
            "- error_type: 오류유형\n"
            "- rationale_bullets: 평가근거(불릿 단위 리스트)\n"
            "- fix_code: 개선방안 코드(회색 박스안에 보이는 모든 코드 추출)\n"
            "- error_code_text: 문제점 코드(회색 박스안에 보이는 모든 코드 추출), 없으면 null\n\n"
            "출력은 반드시 JSON 객체 1개만(아래 템플릿과 동일한 키를 사용):\n"
            "{\"inspection_item\":null,\"error_type\":null,\"rationale_bullets\":[],\"fix_code\":null,\"error_code_text\":null}\n"
            "주의: inspection_item/error_type/fix_code는 문자열 또는 null, rationale_bullets는 문자열 배열, error_code_text는 문자열 또는 null."

        )
        input_payload = [{
            "role": "user",
            "content": [
                {"type": "input_text", "text": prompt},
                {"type": "input_image", "image_url": _data_url_from_image(masked_page_path)},
            ],
        }]

        parsed = self._call_parse_or_fallback(schema_model=SlideFields, input_payload=input_payload)

        if isinstance(parsed, SlideFields):
            parsed.rationale_bullets = [b.strip() for b in parsed.rationale_bullets if b and b.strip()]
            parsed.inspection_item = parsed.inspection_item.strip() if parsed.inspection_item else None
            parsed.error_type = parsed.error_type.strip() if parsed.error_type else None
            parsed.fix_code = parsed.fix_code.strip() if parsed.fix_code else None
            parsed.error_code_text = parsed.error_code_text.strip() if parsed.error_code_text else None
            return parsed

        d = parsed if isinstance(parsed, dict) else {}
        bullets = d.get("rationale_bullets") or []
        if isinstance(bullets, str):
            bullets = [b.strip() for b in re.split(r"[\n\r]+", bullets) if b.strip()]
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        bullets = [str(b).strip() for b in bullets if str(b).strip()]

        return SlideFields(
            inspection_item=(str(d.get("inspection_item")).strip() if d.get("inspection_item") else None),
            error_type=(str(d.get("error_type")).strip() if d.get("error_type") else None),
            rationale_bullets=bullets,
            fix_code=(str(d.get("fix_code")).strip() if d.get("fix_code") else None),
            error_code_text=(str(d.get("error_code_text")).strip() if d.get("error_code_text") else None),
        )


# -----------------------------
# PPTX: 이미지 추출
# -----------------------------
def _iter_picture_shapes(shapes) -> Any:
    for shape in shapes:
        try:
            if MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_picture_shapes(shape.shapes)
            elif MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape
        except Exception:
            continue


def extract_all_images_from_pptx(pptx_path: Path, out_dir: Path):
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다. PPTX 입력을 처리하려면 python-pptx를 설치하세요.")

    prs = Presentation(str(pptx_path))
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)

    img_root = out_dir / "extracted_images" / "pptx"
    ensure_dir(img_root)

    assets: List[ImageAsset] = []
    temp_emu: List[Tuple[int, int, int, int, int, int]] = []  # (si, pi, left, top, w, h)

    slides = list(prs.slides)
    stage("Extract images from PPTX")
    for si, slide in enumerate(pbar(slides, total=len(slides), desc="Extract images (PPTX)", unit="slide")):

        slide_dir = img_root / f"slide_{si:03d}"
        ensure_dir(slide_dir)

        pi = 0
        for shape in _iter_picture_shapes(slide.shapes):
            try:
                img = shape.image
                blob = img.blob
                ext = img.ext or "png"
                h16 = sha1_bytes(blob)[:16]

                left, top, w, h_emu = int(shape.left), int(shape.top), int(shape.width), int(shape.height)
                p = slide_dir / f"img_{pi:03d}_{h16}.{ext}"
                p.write_bytes(blob)

                assets.append(
                    ImageAsset(
                        page_index=si,
                        image_index=pi,
                        path=p,
                        sha1=sha1_bytes(blob),
                        ext=ext,
                        bbox_px=(0, 0, 1, 1),  # placeholder
                        label=None,
                    )
                )
                temp_emu.append((si, pi, left, top, w, h_emu))
                pi += 1
            except Exception:
                continue

    return assets, (slide_w_emu, slide_h_emu), temp_emu


def convert_pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    soffice = _find_soffice()
    ensure_dir(out_dir)
    stage("Convert PPTX -> PDF (LibreOffice)")
    cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)]
    _run_cmd(cmd)

    pdfs = list(out_dir.glob("*.pdf"))
    if not pdfs:
        raise FileNotFoundError("PPTX->PDF 변환 결과 PDF를 찾지 못했습니다.")
    pdfs.sort(key=lambda p: p.stat().st_mtime, reverse=True)
    return pdfs[0]


def render_pdf_to_pngs(pdf_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[Path]:
    ensure_dir(out_dir)
    doc = fitz.open(str(pdf_path))
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)

    png_paths: List[Path] = []
    stage("Render pages to PNG")
    for i in pbar(range(doc.page_count), total=doc.page_count, desc="Render pages", unit="page"):

        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        out_path = out_dir / f"page_{i:03d}.png"
        pix.save(str(out_path))
        png_paths.append(out_path)

    doc.close()
    return png_paths


def render_pages_from_pptx(pptx_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[Path]:
    render_dir = out_dir / "rendered_pages" / "pptx"
    ensure_dir(render_dir)
    pdf_path = convert_pptx_to_pdf(pptx_path, render_dir)
    png_dir = render_dir / "png"
    return render_pdf_to_pngs(pdf_path, png_dir, dpi=dpi)


def compute_bbox_px_for_pptx_assets(assets: List[ImageAsset], pages_png: List[Path], slide_size_emu: Tuple[int, int], temp_emu: List[Tuple[int, int, int, int, int, int]]) -> None:
    slide_w_emu, slide_h_emu = slide_size_emu

    # 페이지별 png 크기
    page_px_size: Dict[int, Tuple[int, int]] = {}
    for i, p in enumerate(pages_png):
        with Image.open(p) as im:
            page_px_size[i] = (im.width, im.height)

    idx = {(a.page_index, a.image_index): a for a in assets}

    stage("Compute image bounding boxes (PPTX)")
    for si, pi, left, top, w_emu, h_emu in pbar(temp_emu, total=len(temp_emu), desc="Compute bboxes", unit="img"):

        if si not in page_px_size:
            continue
        pw, ph = page_px_size[si]

        x = int(round(left / slide_w_emu * pw))
        y = int(round(top / slide_h_emu * ph))
        w = int(round(w_emu / slide_w_emu * pw))
        h = int(round(h_emu / slide_h_emu * ph))

        x = max(0, min(x, pw - 1))
        y = max(0, min(y, ph - 1))
        w = max(1, min(w, pw - x))
        h = max(1, min(h, ph - y))

        a = idx.get((si, pi))
        if a:
            a.bbox_px = (x, y, w, h)



# -----------------------------
# PPTX: 메모/주석(노트/댓글) 추출
# -----------------------------
def _normalize_memo_text(s: str) -> str:
    """메모 텍스트 정규화(줄바꿈 통일 + 앞뒤 공백/빈줄 제거)."""
    s = (s or "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [ln.rstrip() for ln in s.split("\n")]
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    return "\n".join(lines).strip()


def _extract_text_runs(elem: ET.Element) -> str:
    """OOXML 텍스트 런(<a:t>)을 순서대로 결합."""
    parts: List[str] = []
    for t in elem.iter():
        tag = t.tag.split("}")[-1]  # namespace 제거
        if tag == "t" and t.text:
            parts.append(t.text)
    return "".join(parts)


def _extract_modern_comments_from_xml_bytes(xml_bytes: bytes) -> List[str]:
    """ppt/comments/modernComment_*.xml에서 댓글 텍스트만 추출."""
    try:
        root = ET.fromstring(xml_bytes)
    except Exception:
        return []

    out: List[str] = []
    for cm in list(root):
        if cm.tag.split("}")[-1] != "cm":
            continue

        tx = ""
        for child in list(cm):
            if child.tag.split("}")[-1] == "txBody":
                tx = _extract_text_runs(child)
                break
        if not tx:
            tx = _extract_text_runs(cm)

        tx = _normalize_memo_text(tx)
        if tx:
            out.append(tx)
    return out


def _resolve_pptx_comments_target(slide_number_1based: int, target: str) -> str:
    """slideN.xml.rels의 Target(상대경로)을 ZIP 내부 경로로 정규화."""
    # 관계(Relationship)는 slideN.xml 기준 상대경로로 표현됨.
    slide_xml = f"ppt/slides/slide{slide_number_1based}.xml"
    base_dir = posixpath.dirname(slide_xml)
    return posixpath.normpath(posixpath.join(base_dir, target))


def extract_memos_from_pptx(pptx_path: Path) -> Dict[int, str]:
    """
    PPTX에서 메모/주석 텍스트를 페이지(슬라이드) 인덱스(0-based)로 매핑해 반환.

    추출 대상:
    - 발표자 노트(Speaker Notes): python-pptx
    - Modern Comments: ppt/slides/_rels/slideN.xml.rels -> ppt/comments/modernComment_*.xml
    """
    memos_by_slide: Dict[int, List[str]] = {}

    # 1) Speaker Notes
    if Presentation is not None:
        try:
            prs = Presentation(str(pptx_path))
            for si, slide in enumerate(prs.slides):
                try:
                    ns = slide.notes_slide
                    tf = getattr(ns, "notes_text_frame", None) if ns else None
                    txt = _normalize_memo_text(tf.text if tf else "")
                    if txt:
                        memos_by_slide.setdefault(si, []).append(txt)
                except Exception:
                    continue
        except Exception:
            pass

    # 2) Modern Comments (ZIP XML)
    try:
        with zipfile.ZipFile(pptx_path) as z:
            for rels_name in z.namelist():
                m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$", rels_name)
                if not m:
                    continue

                slide_number = int(m.group(1))          # 1-based
                slide_index = slide_number - 1          # 0-based

                try:
                    rel_root = ET.fromstring(z.read(rels_name))
                except Exception:
                    continue

                rel_tag = "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
                for rel in rel_root.findall(rel_tag):
                    typ = (rel.get("Type") or "").lower()
                    if "comments" not in typ:
                        continue

                    target = rel.get("Target") or ""
                    if not target:
                        continue

                    target_path = _resolve_pptx_comments_target(slide_number, target)
                    if target_path not in z.namelist():
                        continue

                    comment_texts = _extract_modern_comments_from_xml_bytes(z.read(target_path))
                    for t in comment_texts:
                        if t:
                            memos_by_slide.setdefault(slide_index, []).append(t)
    except Exception:
        pass

    # 3) 정리/중복제거 + 단일 문자열로 합치기
    out: Dict[int, str] = {}
    for si, items in memos_by_slide.items():
        cleaned: List[str] = []
        seen: set = set()
        for t in items:
            t = _normalize_memo_text(t)
            if not t or t in seen:
                continue
            seen.add(t)
            cleaned.append(t)
        if cleaned:
            out[si] = cleaned[0] if len(cleaned) == 1 else "\n\n---\n\n".join(cleaned)

    return out


# -----------------------------
# PDF: 이미지 추출 + bbox(px)
# -----------------------------
def render_pages_from_pdf(pdf_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[Path]:
    render_dir = out_dir / "rendered_pages" / "pdf" / "png"
    return render_pdf_to_pngs(pdf_path, render_dir, dpi=dpi)


def extract_all_images_from_pdf(pdf_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[ImageAsset]:
    img_root = out_dir / "extracted_images" / "pdf"
    ensure_dir(img_root)

    doc = fitz.open(str(pdf_path))
    zoom = dpi / 72.0

    assets: List[ImageAsset] = []
    stage("Extract images from PDF")
    for page_i in pbar(range(doc.page_count), total=doc.page_count, desc="Extract images (PDF)", unit="page"):

        page = doc.load_page(page_i)
        page_dir = img_root / f"page_{page_i:03d}"
        ensure_dir(page_dir)

        images = page.get_images(full=True)
        img_idx = 0
        for info in images:
            xref = info[0]
            rects = page.get_image_rects(xref)
            if not rects:
                continue

            extracted = doc.extract_image(xref)
            img_bytes = extracted.get("image", b"")
            ext = extracted.get("ext", "png")
            if not img_bytes:
                continue

            h16 = sha1_bytes(img_bytes)[:16]
            img_path = page_dir / f"img_{img_idx:03d}_{h16}.{ext}"
            if not img_path.exists():
                img_path.write_bytes(img_bytes)

            for r in rects:
                x0 = int(round(r.x0 * zoom))
                y0 = int(round(r.y0 * zoom))
                x1 = int(round(r.x1 * zoom))
                y1 = int(round(r.y1 * zoom))
                w = max(1, x1 - x0)
                h = max(1, y1 - y0)

                assets.append(
                    ImageAsset(
                        page_index=page_i,
                        image_index=img_idx,
                        path=img_path,
                        sha1=sha1_bytes(img_bytes),
                        ext=ext,
                        bbox_px=(x0, y0, w, h),
                        label=None,
                    )
                )
            img_idx += 1

    doc.close()
    return assets


# -----------------------------
# 마스킹
# -----------------------------
def mask_page_images(page_png_path: Path, assets_in_page: List[ImageAsset], out_path: Path, pad: int = 2) -> None:
    with Image.open(page_png_path) as im:
        im = im.convert("RGB")
        draw = ImageDraw.Draw(im)
        for a in assets_in_page:
            x, y, w, h = a.bbox_px
            x0 = max(0, x - pad)
            y0 = max(0, y - pad)
            x1 = min(im.width, x + w + pad)
            y1 = min(im.height, y + h + pad)
            draw.rectangle([x0, y0, x1, y1], fill=(255, 255, 255))
        ensure_dir(out_path.parent)
        im.save(out_path)


# -----------------------------
# 공통 파이프라인
# -----------------------------
def run_pipeline(pages_png: List[Path], assets: List[ImageAsset], out_dir: Path, vlm: OpenAIVLMBackend, memos_by_page: Optional[Dict[int, str]] = None) -> Path:
    # 1) 이미지 분류
    stage("Classify images with VLM")
    for a in pbar(assets, total=len(assets), desc="Classify images", unit="img"):

        a.label = vlm.classify_image(a.path)

    # 2) 페이지별 그룹
    assets_by_page: Dict[int, List[ImageAsset]] = {}
    for a in assets:
        assets_by_page.setdefault(a.page_index, []).append(a)

    masked_dir = out_dir / "masked_pages"
    ensure_dir(masked_dir)

    records: List[PageRecord] = []
    stage("Mask pages & extract fields (VLM) / extract code images")
    for page_i, page_png in enumerate(pbar(pages_png, total=len(pages_png), desc="Process pages", unit="page")):

        page_assets = assets_by_page.get(page_i, [])

        # 최종 산출물: other 제외
        error_region_imgs = [a.path for a in page_assets if a.label == "error_region"]
        error_code_imgs = [a.path for a in page_assets if a.label == "error_code"]

        # 텍스트 추출 오염 방지: "모든 이미지 영역" 마스킹
        masked_path = masked_dir / f"page_{page_i:03d}_masked.png"
        mask_page_images(page_png, page_assets, masked_path)

        # 텍스트 필드 추출
        try:
            fields = vlm.extract_fields_from_masked_page(masked_path)
        except Exception as e:
            print(f"[WARN] extract_fields_from_masked_page failed: page={page_i:03d} err={type(e).__name__}: {e}")
            fields = SlideFields()

        # 오류코드 이미지 -> 텍스트 변환(필수)
        code_texts_from_images: List[str] = []
        for img_path in pbar(
            error_code_imgs,
            total=len(error_code_imgs),
            desc=f"Extract code p{page_i:03d}",
            unit="img",
            leave=False,
        ):
            try:
                txt = vlm.extract_code_from_image(img_path)
            except Exception as e:
                print(f"[WARN] extract_code_from_image failed: page={page_i:03d} img={img_path} err={type(e).__name__}: {e}")
                txt = ""
            if txt:
                code_texts_from_images.append(txt)

        error_code_text_from_page = fields.error_code_text.strip() if fields.error_code_text else None

        # 오류코드 최종 우선순위
        if error_code_text_from_page:
            final_error_code_text = error_code_text_from_page
        elif code_texts_from_images:
            final_error_code_text = code_texts_from_images[0] if len(code_texts_from_images) == 1                 else "\n\n/* --- error_code_image_split --- */\n\n".join(code_texts_from_images)
        else:
            final_error_code_text = None

        memo_text = (memos_by_page.get(page_i) if memos_by_page else None)

        records.append(
            PageRecord(
                page_index=page_i,
                page_image_path=page_png,
                page_masked_path=masked_path,
                error_region_images=error_region_imgs,
                error_code_images=error_code_imgs,
                inspection_item=(fields.inspection_item.strip() if fields.inspection_item else None),
                error_type=(fields.error_type.strip() if fields.error_type else None),
                rationale_bullets=[b.strip() for b in fields.rationale_bullets if b and b.strip()],
                fix_code=(fields.fix_code.strip() if fields.fix_code else None),
                error_code_text=final_error_code_text,
                error_code_text_from_page=error_code_text_from_page,
                error_code_texts_from_images=code_texts_from_images,
                memo=memo_text,
            )
        )

    out_jsonl = out_dir / "records.jsonl"
    stage("Write JSONL")
    with out_jsonl.open("w", encoding="utf-8") as f:
        for r in pbar(records, total=len(records), desc="Write records.jsonl", unit="page"):

            obj = dataclasses.asdict(r)
            obj.pop("page_image_path", None)
            obj.pop("page_masked_path", None)
            obj["error_region_images"] = [str(p) for p in r.error_region_images]
            obj["error_code_images"] = [str(p) for p in r.error_code_images]
            f.write(json.dumps(obj, ensure_ascii=False) + "\n")

    return out_jsonl


def preprocess_input(input_path: Path, out_dir: Path, vlm: OpenAIVLMBackend, dpi: int = DEFAULT_DPI) -> Path:
    ensure_dir(out_dir)
    suffix = input_path.suffix.lower()

    if suffix == ".pptx":
        assets, slide_size_emu, temp_emu = extract_all_images_from_pptx(input_path, out_dir)
        pages_png = render_pages_from_pptx(input_path, out_dir, dpi=dpi)
        compute_bbox_px_for_pptx_assets(assets, pages_png, slide_size_emu, temp_emu)
        memos_by_page = extract_memos_from_pptx(input_path)
        return run_pipeline(pages_png, assets, out_dir, vlm, memos_by_page=memos_by_page)

    if suffix == ".pdf":
        pages_png = render_pages_from_pdf(input_path, out_dir, dpi=dpi)
        assets = extract_all_images_from_pdf(input_path, out_dir, dpi=dpi)
        return run_pipeline(pages_png, assets, out_dir, vlm, memos_by_page=None)


    raise ValueError(f"지원하지 않는 확장자입니다: {suffix} (pptx/pdf만 지원)")


# -----------------------------
# CLI
# -----------------------------
def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="입력 파일 경로(.pptx 또는 .pdf)")
    ap.add_argument("--out", required=True, help="출력 디렉터리")
    ap.add_argument("--dpi", type=int, default=DEFAULT_DPI, help=f"렌더링 DPI (기본 {DEFAULT_DPI})")
    ap.add_argument("--model", default=DEFAULT_MODEL, help=f"VLM 모델 (기본 {DEFAULT_MODEL})")
    ap.add_argument("--temperature", type=float, default=DEFAULT_TEMPERATURE, help=f"temperature (기본 {DEFAULT_TEMPERATURE})")
    ap.add_argument("--api-key", default="", help="(선택) OPENAI_API_KEY를 인자로 직접 전달. 미지정 시 환경변수 사용")
    ap.add_argument("--no-progress", action="store_true", help="진행바(tqdm) 비활성화")
    args = ap.parse_args()

    global _PROGRESS_ENABLED
    _PROGRESS_ENABLED = not args.no_progress

    input_path = Path(args.input).expanduser().resolve()
    out_dir = Path(args.out).expanduser().resolve()

    if not input_path.exists():
        raise FileNotFoundError(input_path)

    api_key = args.api_key.strip() or os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY가 필요합니다. 환경변수로 설정하거나 --api-key로 전달하세요.")

    vlm = OpenAIVLMBackend(model=args.model, temperature=args.temperature, api_key=api_key)
    out_jsonl = preprocess_input(input_path, out_dir, vlm, dpi=args.dpi)
    print(f"OK: {out_jsonl}")


if __name__ == "__main__":
    main()
