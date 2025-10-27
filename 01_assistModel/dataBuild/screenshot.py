# screenshot.py
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# python screenshot.py --pptx sample.pptx --metadata metadata.jsonl --width 1366 --scale 1.0 --timeout 60  --scroll-step 1200 --idle-ms 1200 --include-header

import argparse, pathlib, sys, time, subprocess, json, re, os, io
from pathlib import Path
from datetime import datetime
from urllib.parse import urlparse
from typing import Optional, List, Dict, Set
from collections import deque

from playwright.sync_api import sync_playwright, TimeoutError as PWTimeout
from pptx import Presentation

# Pillow는 헤더+본문 합성 옵션에서만 필요
try:
    from PIL import Image
except Exception:
    Image = None

# ================= 기본 설정 =================
DEFAULT_HIDE_SELECTORS = [
    "nav", ".gnb", ".cookie", ".cookies",
    ".banner", ".popup", ".pop", ".floating", ".float",
    ".chat", ".chatbot", ".subscribe", ".sticky",
    ".toolbar", ".footer-quick", ".video-overlay"
]
HEADER_CANDIDATES = ['header', '.header', '#header', '[role="banner"]']

SLIDE_IDX_RE = re.compile(r"(?:^|/|\\)전체영역/(\d+)(?:_\d+)?\.png$", re.IGNORECASE)
URL_RE       = re.compile(r"https?://[^\s<>\"]+", re.IGNORECASE)

# ================ visited 캐시 =================
VISITED_Q: deque[str] = deque(maxlen=256)
VISITED_MAP: Dict[str, bytes] = {}

def canon_url(url: str) -> str:
    p = urlparse(url.strip())
    host = (p.netloc or "").lower()
    scheme = (p.scheme or "").lower()
    path = (p.path or "/").rstrip("/") or "/"
    qs = ("?" + p.query) if p.query else ""
    return f"{scheme}://{host}{path}{qs}"

def cache_get(url: str) -> Optional[bytes]:
    key = canon_url(url)
    data = VISITED_MAP.get(key)
    if data is not None:
        try: VISITED_Q.remove(key)
        except ValueError: pass
        VISITED_Q.append(key)
    return data

def cache_put(url: str, data: bytes, *, maxlen: Optional[int] = None):
    global VISITED_Q
    if maxlen is not None and VISITED_Q.maxlen != maxlen:
        VISITED_Q = deque(list(VISITED_Q), maxlen=maxlen)
    key = canon_url(url)
    if key not in VISITED_MAP and len(VISITED_Q) == VISITED_Q.maxlen:
        old = VISITED_Q.popleft()
        VISITED_MAP.pop(old, None)
    else:
        try: VISITED_Q.remove(key)
        except ValueError: pass
    VISITED_Q.append(key)
    VISITED_MAP[key] = data

# ================ 유틸 ================
def quiet_install_chromium():
    try:
        subprocess.run(
            [sys.executable, "-m", "playwright", "install", "chromium", "--with-deps"],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )
    except Exception:
        pass

def add_stability_styles(page):
    css = """
    * { animation: none !important; transition: none !important; }
    html, body { scroll-behavior: auto !important; }
    """
    try: page.add_style_tag(content=css)
    except Exception: pass

def header_exists(page) -> bool:
    try:
        return page.evaluate(r"""(cands) => cands.some(sel => document.querySelector(sel))""",
                             HEADER_CANDIDATES)
    except Exception:
        return False

def mark_and_hide_header(page, candidates: List[str]):
    page.evaluate(
        r"""(sels) => {
            for (const sel of sels) {
                document.querySelectorAll(sel).forEach(el => {
                    el.setAttribute('data-fullshot-header','1');
                    el.style.setProperty('display','none','important');
                });
            }
        }""",
        HEADER_CANDIDATES,
    )

def show_header_static(page):
    page.evaluate(r"""
    () => {
      document.querySelectorAll('[data-fullshot-header]').forEach(el => {
        el.style.removeProperty('display');
        el.style.setProperty('position','static','important');
        el.style.setProperty('top','auto','important');
        el.style.setProperty('transform','none','important');
        el.querySelectorAll('*').forEach(ch => {
            const cs = getComputedStyle(ch);
            if (cs.position === 'sticky') {
                ch.style.setProperty('position','static','important');
                ch.style.setProperty('top','auto','important');
            }
        });
      });
    }""")

def hide_fixed_and_selectors(page, hide_selectors: List[str], max_h=220):
    try:
        page.evaluate(
            r"""(args) => {
                const sels = args.sels, maxH = args.maxH;
                for (const sel of sels) {
                    document.querySelectorAll(sel).forEach(el => {
                        if (el.closest('[data-fullshot-header]')) return;
                        el.style.setProperty('display','none','important');
                    });
                }
                const all = Array.from(document.querySelectorAll('body *')).slice(0, 5000);
                for (const el of all) {
                    if (el.closest('[data-fullshot-header]')) continue;
                    const cs = getComputedStyle(el);
                    if ((cs.position === 'fixed' || cs.position === 'sticky')) {
                        const r = el.getBoundingClientRect();
                        if (r.height > 0 && r.height <= maxH) {
                            el.style.setProperty('display','none','important');
                        }
                    }
                }
            }""",
            {"sels": hide_selectors, "maxH": int(max_h)}
        )
    except Exception:
        pass

def compose_header_and_body_bytes(header_bytes: bytes, body_bytes: bytes,
                                  out_fmt: str, quality: Optional[int]) -> bytes:
    if Image is None:
        raise RuntimeError("Pillow가 필요합니다.  pip install pillow")
    h = Image.open(io.BytesIO(header_bytes)).convert("RGB")
    b = Image.open(io.BytesIO(body_bytes)).convert("RGB")
    w = min(h.width, b.width)
    if h.width != w: h = h.crop((0,0,w,h.height))
    if b.width != w: b = b.crop((0,0,w,b.height))
    merged = Image.new("RGB", (w, h.height + b.height), (255,255,255))
    merged.paste(h, (0,0))
    merged.paste(b, (0,h.height))
    buf = io.BytesIO()
    if out_fmt.lower() in ("jpg","jpeg"):
        q = 90 if quality is None else max(1, min(100, quality))
        merged.save(buf, format="JPEG", quality=q)
    else:
        merged.save(buf, format="PNG")
    return buf.getvalue()

# ============= 빠른 캡처(페이지 재사용) =============
def capture_bytes_with_page(page, url, *, fmt, quality, include_header,
                            hide_selectors, scroll_mode, step_px, idle_ms) -> bytes:
    add_stability_styles(page)

    # 이동 + 충분 대기
    page.goto(url, wait_until="domcontentloaded")
    try: page.wait_for_load_state("networkidle", timeout=15000)
    except Exception: pass

    # 스크롤: fast(기본) 또는 step
    if scroll_mode == "step":
        last_h = 0; stable = 0
        for _ in range(500):
            page.evaluate("(s)=>window.scrollBy(0,s)", step_px)
            page.wait_for_timeout(idle_ms)
            try:
                h = page.evaluate("() => document.documentElement.scrollHeight")
            except Exception:
                break
            if h == last_h:
                stable += 1
                if stable >= 2: break
            else:
                stable = 0; last_h = h
        page.evaluate("() => window.scrollTo(0, document.documentElement.scrollHeight)")
        page.wait_for_timeout(idle_ms)
    else:
        for _ in range(3):
            page.evaluate("() => window.scrollTo(0, document.documentElement.scrollHeight)")
            page.wait_for_timeout(400)

    # 방해요소/헤더 처리
    if include_header and header_exists(page):
        mark_and_hide_header(page, HEADER_CANDIDATES)
    if hide_selectors:
        hide_fixed_and_selectors(page, hide_selectors, max_h=220)

    if include_header:
        # 본문 먼저
        body = page.screenshot(full_page=True, type="png")
        # 헤더만
        show_header_static(page)
        page.evaluate("() => window.scrollTo(0, 0)")
        page.wait_for_timeout(200)
        try:
            header_el = page.locator("[data-fullshot-header]").first
            header = (header_el.screenshot(type="png")
                      if header_el.count() > 0
                      else page.screenshot(full_page=False, type="png"))
            return compose_header_and_body_bytes(header, body, out_fmt=fmt, quality=quality)
        except Exception:
            # 실패 시 전체 1장으로 대체
            return page.screenshot(full_page=True, type=("jpeg" if fmt=="jpeg" else "png"))
    else:
        return page.screenshot(full_page=True, type=("jpeg" if fmt=="jpeg" else "png"))

# ============= PPTX → 슬라이드별 URL =============
def extract_urls_by_slide(pptx_path: str) -> Dict[int, str]:
    if not os.path.exists(pptx_path):
        print(f"[ERR] PPTX not found: {pptx_path}")
        return {}
    prs = Presentation(pptx_path)
    slide_to_url: Dict[int, str] = {}
    for i, slide in enumerate(prs.slides, start=1):
        urls: List[str] = []
        for shp in slide.shapes:
            try:
                link = getattr(getattr(shp, "click_action", None), "hyperlink", None)
                if link and link.address: urls.append(link.address)
            except Exception: pass
            try:
                link2 = getattr(shp, "hyperlink", None)
                if link2 and link2.address: urls.append(link2.address)
            except Exception: pass
            if getattr(shp, "has_text_frame", False):
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        try:
                            if r.hyperlink and r.hyperlink.address:
                                urls.append(r.hyperlink.address)
                        except Exception: pass
        texts=[]
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                for p in shp.text_frame.paragraphs:
                    texts.append("".join(r.text or "" for r in p.runs))
            if getattr(shp, "has_table", False):
                for r in shp.table.rows:
                    for c in r.cells:
                        if c.text: texts.append(c.text)
        blob = "\n".join(texts)
        for m in URL_RE.finditer(blob): urls.append(m.group(0))
        seen: Set[str] = set()
        for u in urls:
            u2 = u.strip("<> ,.;)]}'\"")
            if u2.lower().startswith(("http://","https://")) and u2 not in seen:
                slide_to_url.setdefault(i, u2); seen.add(u2)
    return slide_to_url

# ============= metadata.jsonl 일괄 재캡처 =============
def recapture_from_metadata(pptx_path: str, metadata_path: str,
                            width: int, scale: float, timeout_s: int,
                            include_header: bool, hide_selectors: List[str],
                            scroll_mode: str, step_px: int, idle_ms: int,
                            fmt: str, quality: Optional[int], channel: Optional[str],
                            cache_size: int, block_media: bool):
    print("[DBG] recapture start")
    print("  PPTX :", os.path.abspath(pptx_path))
    print("  META :", os.path.abspath(metadata_path))

    # visited 캐시 크기 반영
    cache_put("about:blank", b"", maxlen=cache_size)
    VISITED_MAP.pop(canon_url("about:blank"), None)
    try: VISITED_Q.remove(canon_url("about:blank"))
    except ValueError: pass

    slide_urls = extract_urls_by_slide(pptx_path)
    if not os.path.exists(metadata_path):
        print(f"[ERR] metadata.jsonl not found: {metadata_path}")
        return

    processed = 0

    with sync_playwright() as p:
        try:
            browser = p.chromium.launch(headless=True) if channel is None \
                      else p.chromium.launch(headless=True, channel=channel)
        except Exception:
            quiet_install_chromium()
            browser = p.chromium.launch(headless=True)

        context = browser.new_context(
            viewport={"width": width, "height": 900},
            device_scale_factor=scale,
            locale="ko-KR",
            timezone_id="Asia/Seoul"
        )

        # 리소스 차단(옵션)
        if block_media:
            BLOCK_EXT = (".mp4",".webm",".mov",".m4v",".avi",".mp3",".ogg",".wav")
            def _router(route):
                url_l = route.request.url.split("?")[0].lower()
                rtype = route.request.resource_type
                if rtype in {"media"} or url_l.endswith(BLOCK_EXT):
                    return route.abort()
                return route.continue_()
            try:
                context.route("**/*", _router)
            except Exception:
                pass

        page = context.new_page()

        with open(metadata_path, "r", encoding="utf-8") as f:
            for ln, line in enumerate(f, start=1):
                s = line.strip()
                if not s: continue
                try:
                    rec = json.loads(s)
                except Exception as e:
                    print(f"[WARN] line {ln} json error: {e}")
                    continue

                targets = [p for p in (rec.get("file_names", []) or []) if "전체영역" in p]
                for out_path in targets:
                    m = SLIDE_IDX_RE.search(out_path)
                    if not m:
                        print(f"[SKIP] l{ln}: slide idx not found in {out_path}")
                        continue
                    idx = int(m.group(1))
                    url = slide_urls.get(idx)
                    if not url:
                        print(f"[SKIP] slide {idx}: url not found in pptx")
                        continue

                    print(f"[GO] slide {idx} -> {url} -> {out_path}")
                    os.makedirs(os.path.dirname(out_path), exist_ok=True)

                    # 캐시 HIT 시 즉시 저장
                    cached = cache_get(url)
                    if cached is not None:
                        Path(out_path).write_bytes(cached)
                        print(f"[HIT] cache used for {url} ({len(cached)} bytes)")
                        processed += 1
                        continue

                    # 실제 캡처
                    try:
                        img_bytes = capture_bytes_with_page(
                            page, url,
                            fmt=fmt, quality=quality,
                            include_header=include_header,
                            hide_selectors=hide_selectors,
                            scroll_mode=scroll_mode, step_px=step_px, idle_ms=idle_ms
                        )
                        Path(out_path).write_bytes(img_bytes)
                        cache_put(url, img_bytes, maxlen=cache_size)
                        processed += 1
                        print(f"[OK] saved: {out_path} ({len(img_bytes)} bytes)")
                    except PWTimeout:
                        print(f"[TIMEOUT] slide {idx} @ {url}")
                    except Exception as e:
                        print(f"[ERR] slide {idx}: {e}")

        try: context.close()
        finally: browser.close()

    print(f"[DONE] processed: {processed}")

# ================== CLI ==================
def main():
    ap = argparse.ArgumentParser(description="metadata.jsonl + PPTX 기반 전체영역 재캡처 (Playwright, 캐시/재사용/무임시파일)")
    ap.add_argument("--pptx", default="sample.pptx")
    ap.add_argument("--metadata", default="metadata.jsonl")
    ap.add_argument("--width", type=int, default=1366)
    ap.add_argument("--scale", type=float, default=1.0)          # 속도 위해 기본 1.0
    ap.add_argument("--timeout", type=int, default=60)
    ap.add_argument("--fmt", choices=["png","jpeg"], default="png")
    ap.add_argument("--quality", type=int)
    ap.add_argument("--channel", choices=["msedge","chrome"])
    ap.add_argument("--include-header", action="store_true")
    ap.add_argument("--hide-selectors", default=",".join(DEFAULT_HIDE_SELECTORS))
    ap.add_argument("--scroll-mode", choices=["fast","step"], default="fast")
    ap.add_argument("--scroll-step", type=int, default=1200)
    ap.add_argument("--idle-ms", type=int, default=1200)
    ap.add_argument("--cache-size", type=int, default=256)
    ap.add_argument("--no-block-media", action="store_true",
                    help="미디어 차단 비활성화 (기본은 차단)")
    args = ap.parse_args()

    hide_sels = [s.strip() for s in (args.hide_selectors or "").split(",") if s.strip()]
    recapture_from_metadata(
        pptx_path=args.pptx, metadata_path=args.metadata,
        width=args.width, scale=args.scale, timeout_s=args.timeout,
        include_header=args.include_header, hide_selectors=hide_sels,
        scroll_mode=args.scroll_mode, step_px=args.scroll_step, idle_ms=args.idle_ms,
        fmt=args.fmt, quality=args.quality, channel=args.channel,
        cache_size=args.cache_size, block_media=not args.no_block_media
    )

if __name__ == "__main__":
    main()
