# -*- coding: utf-8 -*-
import os, re, json, time, base64
from pptx import Presentation
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from openai import OpenAI

# ===============================
# 출력 디렉토리
# ===============================
os.makedirs("./image/전체영역", exist_ok=True)
os.makedirs("./image/오류영역", exist_ok=True)

# ===============================
# Selenium Driver
# ===============================
def build_driver(headless=True):
    opts = Options()
    if headless:
        opts.add_argument("--headless=new")
    opts.add_argument("--disable-gpu")
    opts.add_argument("--no-sandbox")
    opts.add_argument("--disable-dev-shm-usage")
    opts.add_argument("--window-size=1920,1080")
    opts.add_argument("--force-device-scale-factor=1")
    opts.add_argument("--hide-scrollbars")
    opts.add_argument("--blink-settings=imagesEnabled=true")
    driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=opts)
    driver.set_page_load_timeout(30)
    driver.implicitly_wait(5)
    return driver

def _lazyload_scroll(driver, max_steps=15, step_px=1200, sleep=0.25):
    last_y = 0
    for _ in range(max_steps):
        driver.execute_script(f"window.scrollBy(0, {step_px});")
        time.sleep(sleep)
        y = driver.execute_script("return window.scrollY;") or 0
        if y == last_y:
            break
        last_y = y
    driver.execute_script("window.scrollTo(0, 0);")
    time.sleep(0.2)

def fullpage_screenshot(driver, out_path):
    try:
        driver.execute_cdp_cmd("Page.enable", {})
        metrics = driver.execute_cdp_cmd("Page.getLayoutMetrics", {})
        content = metrics.get("contentSize") or {}
        width = max(1, int(content.get("width", 1920)))
        height = max(1, int(content.get("height", 1080)))
        height = min(height, 40000)
        screenshot = driver.execute_cdp_cmd(
            "Page.captureScreenshot",
            {"format": "png", "fromSurface": True,
             "clip": {"x":0,"y":0,"width":float(width),
                      "height":float(height),"scale":1}}
        )
        with open(out_path, "wb") as f:
            f.write(base64.b64decode(screenshot["data"]))
        return True
    except Exception:
        driver.save_screenshot(out_path)
        return False

# ===============================
# 이미지 보관(오류 영역)
# ===============================
def extract_error_images_from_slide(slide, slide_idx):
    saved, first_written = 0, False
    for shape in slide.shapes:
        is_picture = getattr(shape, "shape_type", None) == 13
        has_image = hasattr(shape, "image") and getattr(shape.image, "blob", None)
        if not (is_picture or has_image): 
            continue
        try:
            blob = shape.image.blob
            if not first_written:
                path = f"./image/오류영역/{slide_idx}.png"; first_written = True
            else:
                path = f"./image/오류영역/{slide_idx}_{saved}.png"
            with open(path, "wb") as f: f.write(blob)
            saved += 1
        except Exception:
            continue
    return saved

# ===============================
# AI 기반 PPTX 파싱
# ===============================
from pptx.enum.shapes import MSO_SHAPE_TYPE

def _collect_text_blocks(slide):
    """슬라이드 내 모든 텍스트/표 블록을 그룹 내부까지 재귀 수집.
    - 줄바꿈/들여쓰기 보존 (정규화/클린업 없음)
    - 그룹 내부 좌표는 그룹의 (left, top)을 누적해 절대좌표로 변환
    """
    blocks = []
    idx = 0

    def add_block(x, y, w, h, text, source):
        nonlocal idx
        text = (text or "").rstrip()
        if text:
            blocks.append({"idx": idx, "x": int(x), "y": int(y), "w": int(w), "h": int(h),
                           "text": text, "source": source})
            idx += 1

    def walk(sh, ox=0, oy=0):
        st = getattr(sh, "shape_type", None)

        # 그룹이면 재귀
        if st == MSO_SHAPE_TYPE.GROUP:
            gx = int(getattr(sh, "left", 0))
            gy = int(getattr(sh, "top", 0))
            for sub in sh.shapes:
                walk(sub, ox + gx, oy + gy)
            return

        # 좌표
        try:
            x = ox + int(getattr(sh, "left", 0))
            y = oy + int(getattr(sh, "top", 0))
            w = int(getattr(sh, "width", 0))
            h = int(getattr(sh, "height", 0))
        except Exception:
            x = ox; y = oy; w = h = 0

        # 표
        if getattr(sh, "has_table", False):
            lines = []
            tbl = sh.table
            for r in tbl.rows:
                for c in r.cells:
                    # 표는 셀 줄바꿈 유지 (정규화 금지)
                    t = (c.text or "")
                    if t.strip() != "":
                        lines.append(t)
            if lines:
                add_block(x, y, w, h, "\n".join(lines), "table")
            return

        # 텍스트 프레임
        if getattr(sh, "has_text_frame", False):
            paras = []
            for p in sh.text_frame.paragraphs:
                # run들을 그대로 이어붙이고 문단 단위로 줄바꿈
                run_text = "".join((r.text or "") for r in p.runs)
                # 코드 같은 경우 공백도 의미가 있으므로 strip() 하지 않음
                if run_text != "":
                    paras.append(run_text)
            if paras:
                add_block(x, y, w, h, "\n".join(paras), "text")
            return

        # 다른 유형은 스킵(그림 등)

    # 최상위에서 시작
    for sh in slide.shapes:
        walk(sh, 0, 0)

    return blocks


client = OpenAI(api_key="sk-proj-RrmGBqNXk-bpIvVhaDFVaIBuV7EQeI3RllDD4M6pDKTzOWocvmXtEsTckK79VVlkLTFOdiqYO6T3BlbkFJnkzL5k4YCpEKAw_bjgJVS6wvlJRbayPKxWRy5BGkyyQWaA7ESppOz6n5Wf25q3hBPZneylOQ4A")

def _ai_map_blocks_to_fields(blocks):
    schema = {
        "페이지명": "string",
        "URL": "string",
        "검사항목": "string",
        "오류유형": "string",
        "문제점": "string",
        "문제점 및 개선방안_텍스트": "string",
        "문제점 및 개선방안_코드": "string"
    }
    blocks_str = json.dumps(blocks, ensure_ascii=False)

    prompt = f"""
    너는 PPT 접근성 진단 보고서의 메타데이터를 **보이는 그대로 전사**하여 구조화하는 도우미다.
    아래는 해당 슬라이드에서 추출한 텍스트/표 블록들의 목록이다(각 블록은 idx,x,y,w,h,text,source 포함).

    [블록 목록]
    {blocks_str}

    [반드시 지켜야 할 원칙]
    - 추측/요약/일반화 금지. **슬라이드에 보이는 문자열을 그대로** 적는다.
    - 출력은 아래 스키마의 **키들을 모두 포함**하는 **단일 JSON 객체 1개**만 허용한다.
    - 어떤 항목을 찾지 못하면 **빈 문자열("")**을 넣는다. (null, N/A, 설명문 금지)
    - 공백/줄바꿈/따옴표/기호(<, >, /, =, :) 등 **문자 하나도 수정하지 않는다**.
    - 코드는 줄바꿈/들여쓰기 포함 **원문 그대로** 보존한다.

    [스키마]
    {json.dumps(schema, ensure_ascii=False, indent=2)}

    [필드 매핑 규칙 — 아주 중요]
    1) '검사항목', '오류유형'
    - 표(헤더 테이블)에서 라벨-값 쌍을 찾아 **값만 그대로** 전사한다.
    - 라벨 예: '검사항목', '오류유형'. 대소문자/공백 변형이 있어도 의미가 같으면 해당 행의 값을 사용.
    - 표에서 못 찾으면 다른 블록에서 '검사항목:', '오류유형:' 뒤의 값을 찾되, 없으면 "".

    2) '문제점'  **← 문제점 섹션의 "코드"만**
    - '문제점' 제목 **아래**에 있는 **코드 박스/코드 라인**만 모아 넣는다.
    - 코드 판단 기준(하나라도 맞으면 코드로 간주):
        • <태그> 형태(예: <img ...>),  • HTML 속성(alt=, class=, id=, role=, aria-),  • 세미콜론이 있는 CSS,  • 주석(<!-- ... -->)
    - 절대 **설명문/불릿 텍스트**를 '문제점'에 넣지 않는다. (예: "장식 등 심미적 효과를 위한 이미지에는..." → 코드가 아님)
    - 코드가 여러 줄/여러 박스면 **원문 순서대로 줄바꿈으로 연결**.

    3) '문제점 및 개선방안_텍스트'
    - '문제점 및 개선방안' 영역의 **서술형 문장**만 담는다. (불릿 포함 가능)
    - 코드/태그/주석은 여기 넣지 않는다.

    4) '문제점 및 개선방안_코드'
    - '개선방안' 제목 **아래**에 제시된 **코드만** 넣는다. (HTML/CSS/JS 등)
    - 여러 조각이면 **원문 순서대로 줄바꿈으로 연결**.

    5) '페이지명', 'URL'
    - 표(헤더)에서 동일 라벨의 값을 그대로 전사한다. 없으면 "".
    - URL 문자열은 원문 그대로(스킴/쿼리 등 수정 금지). 라벨과 값을 혼합해 쓰지 않는다.

    이제 위 규칙에 따라 스키마의 키 순서를 유지하며 **JSON 객체 하나**만 출력하라.
    """

    resp = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role":"system","content":"You are a strict JSON transcriber."},
            {"role":"user","content":prompt},
        ],
        temperature=0.0,
    )
    raw = resp.choices[0].message.content.strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        i = raw.find("\n")
        if i != -1: raw = raw[i+1:].strip()
    return json.loads(raw)

def _sanitize_url(u: str) -> str:
    if not u:
        return ""
    u = u.strip().strip("<>").rstrip(" ,.;)]}'\"")
    return u if u.lower().startswith(("http://", "https://")) else ""

def parse_diagnosis_text(slide, slide_width):
    try:
        blocks = _collect_text_blocks(slide)
        out = _ai_map_blocks_to_fields(blocks)   # ← AI만으로 페이지명/URL/검사항목/오류유형/문제점/개선 텍스트/코드 추출
    except Exception:
        out = {
            "페이지명":"", "URL":"", "검사항목":"", "오류유형":"",
            "문제점":"", "문제점 및 개선방안_텍스트":"", "문제점 및 개선방안_코드":""
        }

    # URL은 AI가 준 값만 사용하되, 형식만 가볍게 확인
    out["URL"] = _sanitize_url(out.get("URL", ""))

    # (선택) 다른 필드들도 꼬리 공백만 제거
    for k in ["페이지명","검사항목","오류유형","문제점","문제점 및 개선방안_텍스트","문제점 및 개선방안_코드"]:
        out[k] = (out.get(k,"") or "").rstrip()

    return out


# ===============================
# 추론(근거 작성)
# ===============================
def generate_reasoning(file_names, output_data):
    full_img = file_names[0] if file_names else "없음"
    err_imgs = file_names[1:] if len(file_names) > 1 else []

    prompt = f"""
    당신은 접근성 근거 분석가입니다.  
    다음 평가 입력에 대한 접근성 진단을 진행했을때, 왜 이런 평가 출력이 나왔는지에 대한 근거(rationale)를 작성해야 합니다.
    평가출력을 그대로 반복하지 말고, 불필요한 설명은 피하십시오.  

    [입력 및 출력]  

    - 평가입력:
    1. 전체 페이지 스크린샷: {full_img}  
    2. 오류 영역 스크린샷: {err_imgs}  
    3. 오류 영역 HTML/코드: {output_data.get("문제점","")}  
    
    - 평가출력:  
    1. 검사항목: {output_data.get("검사항목","")}  
    2. 오류유형: {output_data.get("오류유형","")}  
    3. 개선 텍스트: {output_data.get("문제점 및 개선방안_텍스트","")}  
    4. 개선 코드: {output_data.get("문제점 및 개선방안_코드","")}  

    [작성 지침]  
    1. 전체 페이지 스크린샷을 통해 해당 페이지의 목적을 파악하고, 
    2. 페이지 목적을 참고할 때, 오류 영역 스크린샷에 드러난 진단 콘텐츠의 역할을 파악하십시오.  
    3. 이제 오류 영역 HTML/코드까지 함께 고려할 때, 왜 이런 검사항목과 오류유형이 도출되었는지 설명하십시오.

    [근거]  
    """

    resp = client.chat.completions.create(
        model="o4-mini-2025-04-16",
        messages=[
            {"role":"system","content":"You are a web accessibility expert."},
            {"role":"user","content":prompt},
        ],
    )
    return resp.choices[0].message.content.strip()

# ===============================
# 메인 파이프라인
# ===============================
def process_pptx(pptx_path, output_jsonl="metadata.jsonl", debug=False, headless=True):
    prs = Presentation(pptx_path)
    wrote_any = False

    with open(output_jsonl, "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides, start=1):
            diag = parse_diagnosis_text(slide, prs.slide_width)

            base_url = (diag.get("URL") or "").strip()
            if not base_url:
                if debug: print(f"[DEBUG] Slide {i} URL 없음 → 스킵")
                continue

            img_count = extract_error_images_from_slide(slide, i)
            if debug: print(f"[DEBUG] Slide {i} 오류영역 저장: {img_count}개")

            driver = build_driver(headless=headless)
            input_files = []
            try:
                driver.get(base_url)
                time.sleep(2.0)
                _lazyload_scroll(driver)
                full_path = f"./image/전체영역/{i}.png"
                fullpage_screenshot(driver, full_path)
                input_files.append(full_path)
            except Exception as e:
                if debug: print(f"[DEBUG] Slide {i} 스크린샷 실패: {e}")
                driver.quit()
                continue
            finally:
                try: driver.quit()
                except Exception: pass

            if img_count >= 1:
                for k in range(img_count):
                    path = f"./image/오류영역/{i}.png" if k == 0 else f"./image/오류영역/{i}_{k}.png"
                    if os.path.exists(path): input_files.append(path)

            reasoning = "(추론 실패)"
            try:
                reasoning = generate_reasoning(input_files, diag)
            except Exception as e:
                if debug: print(f"[DEBUG] Slide {i} reasoning 실패: {e}")

            output_obj = {
                "추론": reasoning,
                "검사항목": diag.get("검사항목",""),
                "오류유형": diag.get("오류유형",""),
                "문제점": diag.get("문제점",""),
                "문제점 및 개선방안_텍스트": diag.get("문제점 및 개선방안_텍스트",""),
                "문제점 및 개선방안_코드": diag.get("문제점 및 개선방안_코드",""),
            }
            record = {"file_names": input_files, "output": output_obj}
            f.write(json.dumps(record, ensure_ascii=False) + "\n")
            wrote_any = True
            if debug: print(f"[DEBUG] Slide {i} 완료 → 파일 {len(input_files)}개")

    if debug and not wrote_any:
        print("[DEBUG] 처리된 슬라이드가 없음.")

    return True

if __name__ == "__main__":
    pptx_file = r"C:\Doo\SNC_Lab\01_assistModel\dataBuild\sample.pptx"
    process_pptx(pptx_file, output_jsonl="metadata.jsonl", debug=True, headless=True)
    print("처리 완료")
