# -*- coding: utf-8 -*-
"""
기존 코드의 '추론부'와 '메타데이터 추출부'는 **그대로 유지**하고,
'전체 스크린샷'을 **Playwright 기반**으로 교체했습니다.
- 동적 요소(레이지로드/애니메이션/스티키 헤더) 대응
- 헤더 포함 합성 옵션(기본 포함)
- visited deque 캐시(같은 URL 재방문 시 재캡처 없이 저장)
필요 패키지:
    pip install playwright pillow python-pptx webdriver-manager
    python -m playwright install chromium --with-deps
"""
import os, re, json, time, base64, io
from collections import deque
from urllib.parse import urlparse
from typing import Optional, List, Dict, Set, Tuple
from PIL import Image

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# (Selenium 관련 import/함수는 그대로 두지만 더이상 사용하지 않습니다.)
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager

from openai import OpenAI

# ===============================
# 출력 디렉토리
# ===============================
os.makedirs("./image/전체영역", exist_ok=True)
os.makedirs("./image/오류영역", exist_ok=True)

# ===============================
# (구) Selenium Driver (보존: 사용 안 함)
# ===============================
def build_driver(headless=True):
    opts = Options()
    if headless:
        opts.add_argument("--headless=new")
    opts.add_argument("--disable-gpu")
    opts.add_argument("--no-sandbox")
    opts.add_argument("--disable-dev-shm-usage")
    opts.add_argument("--window-size=1920,1080")
    opts.add_argument("--force-device-scale-factor=1")
    opts.add_argument("--hide-scrollbars")
    opts.add_argument("--blink-settings=imagesEnabled=true")
    driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=opts)
    driver.set_page_load_timeout(30)
    driver.implicitly_wait(5)
    return driver

def _lazyload_scroll(driver, max_steps=15, step_px=1200, sleep=0.25):
    last_y = 0
    for _ in range(max_steps):
        driver.execute_script(f"window.scrollBy(0, {step_px});")
        time.sleep(sleep)
        y = driver.execute_script("return window.scrollY;") or 0
        if y == last_y:
            break
        last_y = y
    driver.execute_script("window.scrollTo(0, 0);")
    time.sleep(0.2)

def fullpage_screenshot(driver, out_path):
    try:
        driver.execute_cdp_cmd("Page.enable", {})
        metrics = driver.execute_cdp_cmd("Page.getLayoutMetrics", {})
        content = metrics.get("contentSize") or {}
        width = max(1, int(content.get("width", 1920)))
        height = max(1, int(content.get("height", 1080)))
        height = min(height, 40000)
        screenshot = driver.execute_cdp_cmd(
            "Page.captureScreenshot",
            {"format": "png", "fromSurface": True,
             "clip": {"x":0,"y":0,"width":float(width),
                      "height":float(height),"scale":1}}
        )
        with open(out_path, "wb") as f:
            f.write(base64.b64decode(screenshot["data"]))
        return True
    except Exception:
        driver.save_screenshot(out_path)
        return False

# ===============================
# 이미지 보관(오류 영역)  ← (원본 유지)
# ===============================
def extract_error_images_from_slide(slide, slide_idx):
    saved, first_written = 0, False
    for shape in slide.shapes:
        is_picture = getattr(shape, "shape_type", None) == 13
        has_image = hasattr(shape, "image") and getattr(shape.image, "blob", None)
        if not (is_picture or has_image): 
            continue
        try:
            blob = shape.image.blob
            if not first_written:
                path = f"./image/오류영역/{slide_idx}.png"; first_written = True
            else:
                path = f"./image/오류영역/{slide_idx}_{saved}.png"
            with open(path, "wb") as f: f.write(blob)
            saved += 1
        except Exception:
            continue
    return saved

# ===============================
# AI 기반 PPTX 파싱  ← (원본 유지)
# ===============================
def _collect_text_blocks(slide):
    """슬라이드 내 모든 텍스트/표 블록을 그룹 내부까지 재귀 수집.
    - 줄바꿈/들여쓰기 보존 (정규화/클린업 없음)
    - 그룹 내부 좌표는 그룹의 (left, top)을 누적해 절대좌표로 변환
    """
    blocks = []
    idx = 0

    def add_block(x, y, w, h, text, source):
        nonlocal idx
        text = (text or "").rstrip()
        if text:
            blocks.append({"idx": idx, "x": int(x), "y": int(y), "w": int(w), "h": int(h),
                           "text": text, "source": source})
            idx += 1

    def walk(sh, ox=0, oy=0):
        st = getattr(sh, "shape_type", None)

        # 그룹이면 재귀
        if st == MSO_SHAPE_TYPE.GROUP:
            gx = int(getattr(sh, "left", 0))
            gy = int(getattr(sh, "top", 0))
            for sub in sh.shapes:
                walk(sub, ox + gx, oy + gy)
            return

        # 좌표
        try:
            x = ox + int(getattr(sh, "left", 0))
            y = oy + int(getattr(sh, "top", 0))
            w = int(getattr(sh, "width", 0))
            h = int(getattr(sh, "height", 0))
        except Exception:
            x = ox; y = oy; w = h = 0

        # 표
        if getattr(sh, "has_table", False):
            lines = []
            tbl = sh.table
            for r in tbl.rows:
                for c in r.cells:
                    # 표는 셀 줄바꿈 유지 (정규화 금지)
                    t = (c.text or "")
                    if t.strip() != "":
                        lines.append(t)
            if lines:
                add_block(x, y, w, h, "\n".join(lines), "table")
            return

        # 텍스트 프레임
        if getattr(sh, "has_text_frame", False):
            paras = []
            for p in sh.text_frame.paragraphs:
                # run들을 그대로 이어붙이고 문단 단위로 줄바꿈
                run_text = "".join((r.text or "") for r in p.runs)
                # 코드 같은 경우 공백도 의미가 있으므로 strip() 하지 않음
                if run_text != "":
                    paras.append(run_text)
            if paras:
                add_block(x, y, w, h, "\n".join(paras), "text")
            return

        # 다른 유형은 스킵(그림 등)

    # 최상위에서 시작
    for sh in slide.shapes:
        walk(sh, 0, 0)

    return blocks

client = OpenAI(api_key="sk-proj-RrmGBqNXk-bpIvVhaDFVaIBuV7EQeI3RllDD4M6pDKTzOWocvmXtEsTckK79VVlkLTFOdiqYO6T3BlbkFJnkzL5k4YCpEKAw_bjgJVS6wvlJRbayPKxWRy5BGkyyQWaA7ESppOz6n5Wf25q3hBPZneylOQ4A")

def _ai_map_blocks_to_fields(blocks):
    schema = {
        "페이지명": "string",
        "URL": "string",
        "검사항목": "string",
        "오류유형": "string",
        "문제점": "string",
        "문제점 및 개선방안_텍스트": "string",
        "문제점 및 개선방안_코드": "string"
    }
    blocks_str = json.dumps(blocks, ensure_ascii=False)

    prompt = f"""
    너는 PPT 접근성 진단 보고서의 메타데이터를 **보이는 그대로 전사**하여 구조화하는 도우미다.
    아래는 해당 슬라이드에서 추출한 텍스트/표 블록들의 목록이다(각 블록은 idx,x,y,w,h,text,source 포함).

    [블록 목록]
    {blocks_str}

    [반드시 지켜야 할 원칙]
    - 추측/요약/일반화 금지. **슬라이드에 보이는 문자열을 그대로** 적는다.
    - 출력은 아래 스키마의 **키들을 모두 포함**하는 **단일 JSON 객체 1개**만 허용한다.
    - 어떤 항목을 찾지 못하면 **빈 문자열("")**을 넣는다. (null, N/A, 설명문 금지)
    - 공백/줄바꿈/따옴표/기호(<, >, /, =, :) 등 **문자 하나도 수정하지 않는다**.
    - 코드는 줄바꿈/들여쓰기 포함 **원문 그대로** 보존한다.

    [스키마]
    {json.dumps(schema, ensure_ascii=False, indent=2)}

    [필드 매핑 규칙 — 아주 중요]
    1) '검사항목', '오류유형'
    - 표(헤더 테이블)에서 라벨-값 쌍을 찾아 **값만 그대로** 전사한다.
    - 라벨 예: '검사항목', '오류유형'. 대소문자/공백 변형이 있어도 의미가 같으면 해당 행의 값을 사용.
    - 표에서 못 찾으면 다른 블록에서 '검사항목:', '오류유형:' 뒤의 값을 찾되, 없으면 "".

    2) '문제점'  **← 문제점 섹션의 "코드"만**
    - '문제점' 제목 **아래**에 있는 **코드 박스/코드 라인**만 모아 넣는다.
    - 코드 판단 기준(하나라도 맞으면 코드로 간주):
        • <태그> 형태(예: <img ...>),  • HTML 속성(alt=, class=, id=, role=, aria-),  • 세미콜론이 있는 CSS,  • 주석(<!-- ... -->)
    - 절대 **설명문/불릿 텍스트**를 '문제점'에 넣지 않는다. (예: "장식 등 심미적 효과를 위한 이미지에는..." → 코드가 아님)
    - 코드가 여러 줄/여러 박스면 **원문 순서대로 줄바꿈으로 연결**.

    3) '문제점 및 개선방안_텍스트'
    - '문제점 및 개선방안' 영역의 **서술형 문장**만 담는다. (불릿 포함 가능)
    - 코드/태그/주석은 여기 넣지 않는다.

    4) '문제점 및 개선방안_코드'
    - '개선방안' 제목 **아래**에 제시된 **코드만** 넣는다. (HTML/CSS/JS 등)
    - 여러 조각이면 **원문 순서대로 줄바꿈으로 연결**.

    5) '페이지명', 'URL'
    - 표(헤더)에서 동일 라벨의 값을 그대로 전사한다. 없으면 "".
    - URL 문자열은 원문 그대로(스킴/쿼리 등 수정 금지). 라벨과 값을 혼합해 쓰지 않는다.

    이제 위 규칙에 따라 스키마의 키 순서를 유지하며 **JSON 객체 하나**만 출력하라.
    """

    resp = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role":"system","content":"You are a strict JSON transcriber."},
            {"role":"user","content":prompt},
        ],
        temperature=0.0,
    )
    raw = resp.choices[0].message.content.strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        i = raw.find("\n")
        if i != -1: raw = raw[i+1:].strip()
    return json.loads(raw)

def _sanitize_url(u: str) -> str:
    if not u:
        return ""
    u = u.strip().strip("<>").rstrip(" ,.;)]}'\"")
    return u if u.lower().startswith(("http://", "https://")) else ""

def parse_diagnosis_text(slide, slide_width):
    try:
        blocks = _collect_text_blocks(slide)
        out = _ai_map_blocks_to_fields(blocks)   # ← AI만으로 페이지명/URL/검사항목/오류유형/문제점/개선 텍스트/코드 추출
    except Exception:
        out = {
            "페이지명":"", "URL":"", "검사항목":"", "오류유형":"",
            "문제점":"", "문제점 및 개선방안_텍스트":"", "문제점 및 개선방안_코드":""
        }

    # URL은 AI가 준 값만 사용하되, 형식만 가볍게 확인
    out["URL"] = _sanitize_url(out.get("URL", ""))

    # (선택) 다른 필드들도 꼬리 공백만 제거
    for k in ["페이지명","검사항목","오류유형","문제점","문제점 및 개선방안_텍스트","문제점 및 개선방안_코드"]:
        out[k] = (out.get(k,"") or "").rstrip()

    return out

# ===============================
# 추론(근거 작성)  ← (원본 유지)
# ===============================
def generate_reasoning(file_names, output_data):
    full_img = file_names[0] if file_names else "없음"
    err_imgs = file_names[1:] if len(file_names) > 1 else []

    prompt = f"""
    당신은 접근성 근거 분석가입니다.  
    다음 평가 입력에 대한 접근성 진단을 진행했을때, 왜 이런 평가 출력이 나왔는지에 대한 근거(rationale)를 작성해야 합니다.
    평가출력을 그대로 반복하지 말고, 불필요한 설명은 피하십시오.  

    [입력 및 출력]  

    - 평가입력:
    1. 전체 페이지 스크린샷: {full_img}  
    2. 오류 영역 스크린샷: {err_imgs}  
    3. 오류 영역 HTML/코드: {output_data.get("문제점","")}  
    
    - 평가출력:  
    1. 검사항목: {output_data.get("검사항목","")}  
    2. 오류유형: {output_data.get("오류유형","")}  
    3. 개선 텍스트: {output_data.get("문제점 및 개선방안_텍스트","")}  
    4. 개선 코드: {output_data.get("문제점 및 개선방안_코드","")}  

    [작성 지침]  
    1. 전체 페이지 스크린샷을 통해 해당 페이지의 목적을 파악하고, 
    2. 페이지 목적을 참고할 때, 오류 영역 스크린샷에 드러난 진단 콘텐츠의 역할을 파악하십시오.  
    3. 이제 오류 영역 HTML/코드까지 함께 고려할 때, 왜 이런 [평가 출력]이 도출되었는지 설명하십시오.

    [근거]  
    """

    resp = client.responses.create(
    model="gpt-5",
    input=[
        {"role": "system", "content": "You are a web accessibility expert."},
        {"role": "user", "content": prompt},
    ],
    )

    reasoning = resp.output_text.strip()
    return reasoning

# =====================================================================
# ★★★ Playwright 기반 전체 스크린샷 (동적 로딩/헤더 합성/캐시) — 신규 추가 ★★★
# =====================================================================
from playwright.sync_api import sync_playwright, TimeoutError as PWTimeout

DEFAULT_HIDE_SELECTORS = [
    "nav", ".gnb", ".cookie", ".cookies",
    ".banner", ".popup", ".pop", ".floating", ".float",
    ".chat", ".chatbot", ".subscribe", ".sticky",
    ".toolbar", ".footer-quick", ".video-overlay"
]
HEADER_CANDIDATES = ['header', '.header', '#header', '[role="banner"]']

# visited LRU 캐시
_VISITED_Q: deque[str] = deque(maxlen=256)
_VISITED_MAP: Dict[str, bytes] = {}

def _canon_url(u: str) -> str:
    p = urlparse((u or "").strip())
    scheme = (p.scheme or "").lower()
    host = (p.netloc or "").lower()
    path = (p.path or "/").rstrip("/") or "/"
    qs = ("?" + p.query) if p.query else ""
    return f"{scheme}://{host}{path}{qs}"

def _cache_get(u: str) -> Optional[bytes]:
    k = _canon_url(u)
    b = _VISITED_MAP.get(k)
    if b is not None:
        try: _VISITED_Q.remove(k)
        except ValueError: pass
        _VISITED_Q.append(k)
    return b

def _cache_put(u: str, b: bytes):
    k = _canon_url(u)
    if k not in _VISITED_MAP and len(_VISITED_Q) == _VISITED_Q.maxlen:
        old = _VISITED_Q.popleft()
        _VISITED_MAP.pop(old, None)
    else:
        try: _VISITED_Q.remove(k)
        except ValueError: pass
    _VISITED_Q.append(k)
    _VISITED_MAP[k] = b

def _add_stability_styles(page):
    css = """
    * { animation: none !important; transition: none !important; }
    html, body { scroll-behavior: auto !important; }
    """
    try: page.add_style_tag(content=css)
    except Exception: pass

def _header_exists(page) -> bool:
    try:
        return page.evaluate(
            "(sels)=>sels.some(sel=>document.querySelector(sel))",
            HEADER_CANDIDATES
        )
    except Exception:
        return False

def _mark_and_hide_header(page):
    page.evaluate(
        "(sels)=>{for(const sel of sels){document.querySelectorAll(sel).forEach(el=>{el.setAttribute('data-fullshot-header','1');el.style.setProperty('display','none','important');});}}",
        HEADER_CANDIDATES
    )

def _show_header_static(page):
    page.evaluate("""
    () => {
      document.querySelectorAll('[data-fullshot-header]').forEach(el => {
        el.style.removeProperty('display');
        el.style.setProperty('position','static','important');
        el.style.setProperty('top','auto','important');
        el.style.setProperty('transform','none','important');
        el.querySelectorAll('*').forEach(ch => {
            const cs = getComputedStyle(ch);
            if (cs.position === 'sticky') {
                ch.style.setProperty('position','static','important');
                ch.style.setProperty('top','auto','important');
            }
        });
      });
    }""")

def _hide_fixed_and_selectors(page, hide_selectors: List[str], max_h=220):
    try:
        page.evaluate(
            """(args)=>{
                const sels=args.sels, maxH=args.maxH;
                for (const sel of sels) {
                    document.querySelectorAll(sel).forEach(el=>{
                        if (el.closest('[data-fullshot-header]')) return;
                        el.style.setProperty('display','none','important');
                    });
                }
                const all=[...document.querySelectorAll('body *')].slice(0,5000);
                for(const el of all){
                    if (el.closest('[data-fullshot-header]')) continue;
                    const cs=getComputedStyle(el);
                    if(cs.position==='fixed'||cs.position==='sticky'){
                        const r=el.getBoundingClientRect();
                        if(r.height>0 && r.height<=maxH){
                            el.style.setProperty('display','none','important');
                        }
                    }
                }
            }""",
            {"sels": hide_selectors, "maxH": int(max_h)}
        )
    except Exception:
        pass

def _compose_header_and_body_bytes(header_bytes: bytes, body_bytes: bytes, *, out_fmt: str = "png", quality: Optional[int] = None) -> bytes:
    if Image is None:
        raise RuntimeError("Pillow가 필요합니다.  pip install pillow")
    h = Image.open(io.BytesIO(header_bytes)).convert("RGB")
    b = Image.open(io.BytesIO(body_bytes)).convert("RGB")
    w = min(h.width, b.width)
    if h.width != w: h = h.crop((0,0,w,h.height))
    if b.width != w: b = b.crop((0,0,w,b.height))
    merged = Image.new("RGB", (w, h.height + b.height), (255,255,255))
    merged.paste(h, (0,0)); merged.paste(b, (0,h.height))
    buf = io.BytesIO()
    if out_fmt.lower() in ("jpg","jpeg"):
        q = 90 if quality is None else max(1, min(100, quality))
        merged.save(buf, format="JPEG", quality=q)
    else:
        merged.save(buf, format="PNG")
    return buf.getvalue()

def _capture_with_page(page, url: str, *, width=1366, scale=1.0, timeout_s=60,
                       include_header=True, hide_selectors: Optional[List[str]]=None,
                       scroll_mode="fast", step_px=1200, idle_ms=1200,
                       block_media=True, out_fmt="png", quality: Optional[int]=None) -> bytes:
    hide_selectors = hide_selectors or DEFAULT_HIDE_SELECTORS

    _add_stability_styles(page)

    page.set_viewport_size({"width": width, "height": 900})
    page.goto(url, wait_until="domcontentloaded", timeout=timeout_s*1000)
    try:
        page.wait_for_load_state("networkidle", timeout=timeout_s*1000)
    except PWTimeout:
        pass

    # 스크롤 로딩: fast(3회 바닥 점프) 또는 step
    if scroll_mode == "step":
        last_h = 0; stable = 0
        for _ in range(500):
            page.evaluate("(s)=>window.scrollBy(0,s)", step_px)
            page.wait_for_timeout(idle_ms)
            try:
                h = page.evaluate("()=>document.documentElement.scrollHeight")
            except Exception:
                break
            if h == last_h:
                stable += 1
                if stable >= 2: break
            else:
                stable = 0; last_h = h
        page.evaluate("()=>window.scrollTo(0,document.documentElement.scrollHeight)")
        page.wait_for_timeout(idle_ms)
    else:
        for _ in range(3):
            page.evaluate("()=>window.scrollTo(0,document.documentElement.scrollHeight)")
            page.wait_for_timeout(400)

    # 방해 요소 및 헤더 처리
    if include_header and _header_exists(page):
        _mark_and_hide_header(page)
    if hide_selectors:
        _hide_fixed_and_selectors(page, hide_selectors, max_h=220)

    if include_header:
        # 본문 먼저
        body = page.screenshot(full_page=True, type="png")
        # 헤더만
        _show_header_static(page)
        page.evaluate("()=>window.scrollTo(0,0)")
        page.wait_for_timeout(200)
        try:
            header_el = page.locator("[data-fullshot-header]").first
            header = header_el.screenshot(type="png") if header_el.count() > 0 \
                     else page.screenshot(full_page=False, type="png")
            final = _compose_header_and_body_bytes(header, body, out_fmt=out_fmt, quality=quality)
        except Exception:
            final = page.screenshot(full_page=True, type=("jpeg" if out_fmt=="jpeg" else "png"))
    else:
        final = page.screenshot(full_page=True, type=("jpeg" if out_fmt=="jpeg" else "png"))

    return final

def _open_playwright(width=1366, scale=1.0, block_media=True):
    """브라우저/컨텍스트/페이지 생성 (한 번만 열어 재사용)."""
    pw = sync_playwright().start()
    try:
        browser = pw.chromium.launch(headless=True)
    except Exception:
        browser = pw.chromium.launch(headless=True)
    context = browser.new_context(
        viewport={"width": width, "height": 900},
        device_scale_factor=scale,
        locale="ko-KR",
        timezone_id="Asia/Seoul"
    )
    if block_media:
        BLOCK_EXT = (".mp4",".webm",".mov",".m4v",".avi",".mp3",".ogg",".wav")
        def _router(route):
            url_l = route.request.url.split("?")[0].lower()
            rtype = route.request.resource_type
            if rtype in {"media"} or url_l.endswith(BLOCK_EXT):
                return route.abort()
            return route.continue_()
        try:
            context.route("**/*", _router)
        except Exception:
            pass
    page = context.new_page()
    return pw, browser, context, page

def _close_playwright(pw, browser, context):
    try:
        context.close()
    except Exception:
        pass
    try:
        browser.close()
    except Exception:
        pass
    try:
        pw.stop()
    except Exception:
        pass

# ===============================
# 번호 이어붙이기 유틸 (추가)
# ===============================
def _max_index_in(dir_path: str) -> int:
    max_id = 0
    if not os.path.isdir(dir_path):
        return 0
    for name in os.listdir(dir_path):
        m = re.search(r'^(\d+)\.png$', name, re.IGNORECASE)
        if m:
            try:
                max_id = max(max_id, int(m.group(1)))
            except ValueError:
                pass
    return max_id


# ===============================
# 메인 파이프라인  ← (★ 스크린샷 부분만 Playwright로 교체)
# ===============================
def process_pptx(pptx_path, output_jsonl="metadata.jsonl", debug=False, headless=True):
    prs = Presentation(pptx_path)
    wrote_any = False

    # ✅ 여기서 페이지/브라우저 열기 (한 번 열어 재사용)
    pw, browser, context, page = _open_playwright(width=1366, scale=1.0, block_media=True)

    # ✅ 전체영역 파일 번호 이어붙이기(이미 추가했다면 유지)
    start_full_idx = _max_index_in("./image/전체영역")

    try:
        # ✅ 메타데이터는 append 모드로
        with open(output_jsonl, "a", encoding="utf-8") as f:
            for i, slide in enumerate(prs.slides, start=1):
                diag = parse_diagnosis_text(slide, prs.slide_width)

                base_url = (diag.get("URL") or "").strip()
                if not base_url:
                    if debug: print(f"[DEBUG] Slide {i} URL 없음 → 스킵")
                    continue

                # 전역(이어붙이기) 인덱스
                global_idx = start_full_idx + i

                # 오류영역 저장(전역 번호 사용)
                img_count = extract_error_images_from_slide(slide, global_idx)
                if debug: print(f"[DEBUG] Slide {i} 오류영역 저장: {img_count}개")

                input_files = []
                try:
                    full_path = f"./image/전체영역/{global_idx}.png"

                    # 캐시 확인
                    cached = _cache_get(base_url)
                    if cached is not None:
                        with open(full_path, "wb") as out:
                            out.write(cached)
                        if debug: print(f"[DEBUG] Slide {i} 캐시 HIT → 저장만")
                    else:
                        # ✅ Playwright 페이지 재사용해 전체 캡처
                        img_bytes = _capture_with_page(
                            page, base_url,
                            width=1366, scale=1.0, timeout_s=60,
                            include_header=True, hide_selectors=DEFAULT_HIDE_SELECTORS,
                            scroll_mode="fast", step_px=1200, idle_ms=1200,
                            block_media=True, out_fmt="png", quality=None
                        )
                        with open(full_path, "wb") as out:
                            out.write(img_bytes)
                        _cache_put(base_url, img_bytes)
                        if debug: print(f"[DEBUG] Slide {i} 스크린샷 완료 → {full_path}")

                    input_files.append(full_path)
                except Exception as e:
                    if debug: print(f"[DEBUG] Slide {i} 스크린샷 실패: {e}")
                    continue

                # 오류영역 경로 수집(전역 번호 기준)
                if img_count >= 1:
                    for k in range(img_count):
                        path = f"./image/오류영역/{global_idx}.png" if k == 0 else f"./image/오류영역/{global_idx}_{k}.png"
                        if os.path.exists(path): input_files.append(path)

                # 추론부/메타데이터 기록(원본 유지)
                reasoning = "(추론 실패)"
                try:
                    reasoning = generate_reasoning(input_files, diag)
                except Exception as e:
                    if debug: print(f"[DEBUG] Slide {i} reasoning 실패: {e}")

                output_obj = {
                    "추론": reasoning,
                    "검사항목": diag.get("검사항목",""),
                    "오류유형": diag.get("오류유형",""),
                    "문제점": diag.get("문제점",""),
                    "문제점 및 개선방안_텍스트": diag.get("문제점 및 개선방안_텍스트",""),
                    "문제점 및 개선방안_코드": diag.get("문제점 및 개선방안_코드",""),
                }
                record = {"file_names": input_files, "output": output_obj}
                f.write(json.dumps(record, ensure_ascii=False) + "\n")
                wrote_any = True
                if debug: print(f"[DEBUG] Slide {i} 완료 → 파일 {len(input_files)}개")

        if debug and not wrote_any:
            print("[DEBUG] 처리된 슬라이드가 없음.")
        return True
    finally:
        # ✅ 여기서 닫기
        _close_playwright(pw, browser, context)

# ===============================
# 진입점 (원본 유지)
# ===============================
if __name__ == "__main__":
    pptx_file = r"C:\Doo\SNC_Lab\01_assistModel\dataBuild\sample11_3.pptx"
    process_pptx(pptx_file, output_jsonl="metadata.jsonl", debug=True, headless=True)
    print("처리 완료")
