#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract images and accessibility evaluation metadata from PPTX slides, and
generate a rationale using a reasoning model (e.g., GPT-4o).

Outputs:
- ./train/<pptx_basename>/<image files>.png
- ./train/<pptx_basename>.jsonl (one JSON object per slide that had any data)

Usage:
  python extract_pptx_accessibility.py example.pptx \
    --model gpt-4o-mini --openai_api_key $OPENAI_API_KEY

If you omit the API key or pass --dry-run, the script will skip calling the model
and leave "rationale" empty.
"""

from __future__ import annotations

import os
import re
import json
import argparse
from typing import Dict, List, Tuple, Any, Optional

# Third-party
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Optional: OpenAI client (lazy import only if needed)
def _lazy_openai_client(api_key: str):
    try:
        # Newer SDK style
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        return client, "responses"
    except Exception:
        try:
            import openai
            openai.api_key = api_key
            return openai, "chat.completions"
        except Exception as e:
            raise RuntimeError("OpenAI SDK not installed. Add 'openai' to requirements and install.") from e


PROMPT_TEMPLATE = """You are an accessibility rationale analyst. Given a visual screenshot of a user interface,
and an accessibility evaluation result (including evaluated items, reason and improvement plan),
your task is to generate a rationale that explains why the accessibility issue was correctly identified.

You must refer to both the visual structure and the code implementation, and where applicable,
link the issue to a specific error type from the provided guideline (e.g. WCAG 5.1.1 error types 1-1 to 1-11).
Ensure adherence to these guidelines:

1. Visually interpret the screenshot and identify the function or role of evaluated elements (e.g. decorative, instructional, interactive).
2. Inspect the HTML snippet to determine whether semantic markup, alt, aria-*, or other accessibility features are properly applied.
3. Compare the visual intent and HTML implementation:
  - Is there a meaningful image that lacks a descriptive alt text?
  - Is decorative content improperly exposed?
  - Are users relying on visual-only cues?
4. Reference the appropriate rule from the guideline table that supports the identified issue.
5. Explain how the current implementation fails to meet that guideline and how it impacts accessibility (especially for screen reader users).
6. Do not restate the original evaluation reason or output the proposed fix. Only provide the rationale
7. Do not question me. Just generate rationale

Screenshot (path): {image_path}

Guideline table:
{guide}

Evaluation result:
{eval}

Rationale:
"""


# --- Constants / Aliases ------------------------------------------------------

KOREAN_FIELD_ALIASES = {
    "페이지명": ["페이지명", "페이지 이름", "Page Name"],
    "URL": ["URL", "Url", "url"],
    "검사항목": ["검사항목", "평가항목", "Evaluated Item", "검사 항목"],
    "오류유형": ["오류유형", "오류 유형", "Error Type"],
    "문제점 및 개선 방안(텍스트)": ["문제점 및 개선 방안", "문제점 및 개선 방안(텍스트)", "문제점", "개선 방안 텍스트"],
    "개선 방안(코드)": ["개선 방안(코드)", "개선 방안 코드", "코드", "개선방안(코드)"],
}

FIELD_KEYS_ORDER = [
    "페이지명", "URL", "검사항목", "오류유형",
    "문제점 및 개선 방안(텍스트)", "개선 방안(코드)",
]

LABEL_RE = r"(?:{})(?:\s*[:：]\s*|\s+)"  # label followed by : or space


# --- Helpers ------------------------------------------------------------------

def _norm_text(s: str) -> str:
    return re.sub(r"\s+", " ", s.strip())


def _iter_shapes_recursive(container, parent_offset: Tuple[int, int]=(0, 0)):
    """Yield (shape, abs_left, abs_top) recursively from slide or group, in reading-ish order."""
    items: List[Tuple[Any, int, int]] = []
    for shp in getattr(container, "shapes", []):
        try:
            left = int(getattr(shp, "left", 0)) + parent_offset[0]
            top = int(getattr(shp, "top", 0)) + parent_offset[1]
        except Exception:
            left, top = parent_offset
        items.append((shp, left, top))

    # Sort roughly by top then left for stable reading order
    items.sort(key=lambda t: (t[2], t[1]))

    for shp, left, top in items:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Yield group (in case it has text), then children
            yield shp, left, top
            for inner in _iter_shapes_recursive(shp, (left, top)):
                yield inner
        else:
            yield shp, left, top


def _extract_images_from_slide(slide, out_dir: str, slide_idx: int) -> List[str]:
    """Save images found on the slide (including inside groups) and return saved paths."""
    os.makedirs(out_dir, exist_ok=True)
    saved_paths: List[str] = []
    img_num = 0
    for shp, _, _ in _iter_shapes_recursive(slide):
        try:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shp.image
                ext = (img.ext or "png").lower()
                if not ext.startswith("."):
                    ext = "." + ext
                filename = f"s{slide_idx:03d}_p{img_num:02d}{ext}"
                path = os.path.join(out_dir, filename)
                with open(path, "wb") as f:
                    f.write(img.blob)
                saved_paths.append(path)
                img_num += 1
        except Exception:
            # Skip problematic image without crashing the whole slide
            continue
    return saved_paths


def _extract_text_from_slide(slide, debug_sink: Optional[Any]=None) -> List[str]:
    """Collect text from text frames and tables, recursively, roughly in reading order."""
    chunks: List[str] = []
    for shp, _, _ in _iter_shapes_recursive(slide):
        try:
            if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = shp.table
                for r in tbl.rows:
                    cells = [c.text.strip() for c in r.cells]
                    if any(cells):
                        chunks.append(" | ".join([c for c in cells if c]))
                continue
            if getattr(shp, "has_text_frame", False):
                t = shp.text
                if t and t.strip():
                    for line in t.splitlines():
                        line = line.strip()
                        if line:
                            chunks.append(line)
        except Exception:
            continue

    if debug_sink is not None:
        try:
            debug_sink.write("[RAW TEXT CHUNKS]\n")
            for i, c in enumerate(chunks, 1):
                debug_sink.write(f"{i:03d}: {c}\n")
        except Exception:
            pass
    return chunks


def _parse_fields_from_texts(all_texts: List[str]) -> Dict[str, str]:
    """Parse labeled values from free text + tables robustly (Korean labels).
       추가: 라벨이 없을 때도 값들을 휴리스틱으로 추정해 채운다.
    """
    fields: Dict[str, str] = {k: "" for k in FIELD_KEYS_ORDER}
    lines = [re.sub(r"\s+", " ", s.strip()) for s in all_texts if s and s.strip()]

    # ---- 1) 라벨 기반 (기존 로직)
    alias_map = {alias: key for key, aliases in KOREAN_FIELD_ALIASES.items() for alias in aliases}
    for ln in lines:
        if " | " in ln:
            parts = [p.strip() for p in ln.split(" | ") if p.strip()]
            if len(parts) >= 2:
                label = parts[0]
                if label in alias_map and not fields[alias_map[label]]:
                    fields[alias_map[label]] = " ".join(parts[1:]).strip()

    for ln in lines:
        for alias, key in alias_map.items():
            m = re.match(rf"^{re.escape(alias)}\s*[:：]?\s*(.+)$", ln)
            if m and not fields[key]:
                fields[key] = m.group(1).strip()

    def collect_block(start_aliases: List[str], stop_keys: List[str]) -> str:
        start_idx = -1
        for i, ln in enumerate(lines):
            if any(ln == a or ln.startswith(a) for a in start_aliases):
                start_idx = i
                break
        if start_idx == -1:
            return ""
        stop_aliases = sum([KOREAN_FIELD_ALIASES[k] for k in stop_keys], [])
        blk = []
        for j in range(start_idx + 1, len(lines)):
            if any(lines[j] == sa or lines[j].startswith(sa) for sa in stop_aliases):
                break
            blk.append(lines[j])
        return "\n".join(blk).strip()

    # 간단 필드가 비어있으면 "다음 줄 값" 휴리스틱
    for key in ["페이지명", "URL", "검사항목", "오류유형"]:
        if not fields[key]:
            aliases = KOREAN_FIELD_ALIASES[key]
            for i, ln in enumerate(lines):
                if any(ln == a or ln.startswith(a) for a in aliases):
                    if i + 1 < len(lines):
                        nxt = lines[i + 1]
                        if nxt and nxt not in alias_map:
                            fields[key] = nxt
                            break

    # 블록 필드
    if not fields["문제점 및 개선 방안(텍스트)"]:
        fields["문제점 및 개선 방안(텍스트)"] = collect_block(
            KOREAN_FIELD_ALIASES["문제점 및 개선 방안(텍스트)"],
            ["개선 방안(코드)", "페이지명", "URL", "검사항목", "오류유형"],
        )
    if not fields["개선 방안(코드)"]:
        fields["개선 방안(코드)"] = collect_block(
            KOREAN_FIELD_ALIASES["개선 방안(코드)"],
            ["문제점 및 개선 방안(텍스트)", "페이지명", "URL", "검사항목", "오류유형"],
        )

    # ---- 2) 라벨이 전혀 없는 슬라이드용 휴리스틱 (여기서부터 추가)
    # 이미 채워졌다면 건너뜀
    all_set = any(fields.values())
    if not all_set:
        # a) URL: 첫 번째 URL
        url_re = re.compile(
            r"(https?://[^\s)]+)",
            re.IGNORECASE,
        )
        urls = []
        for ln in lines:
            urls.extend(url_re.findall(ln))
        if urls and not fields["URL"]:
            fields["URL"] = urls[0]

        # b) 페이지명: 경로 형태(예: "영어_초56 > 커스텀 > 9단원 > 3차시") 또는
        #    첫 줄(제목/헤더)을 제외한 "의미 있어 보이는" 첫 줄
        def looks_like_path(s: str) -> bool:
            return (" > " in s) or ("단원" in s) or ("차시" in s)

        candidate_page = ""
        for ln in lines:
            if looks_like_path(ln) and ln != "진단 결과 및 개선 방안":
                candidate_page = ln
                break
        if not candidate_page:
            # 헤더성 문구를 건너뛰고, URL/가이드라인 설명이 아닌 첫 문장
            skip_phrases = ("진단 결과", "개선 방안", "동일 유형 페이지")
            for ln in lines:
                if any(sp in ln for sp in skip_phrases):
                    continue
                if url_re.search(ln):
                    continue
                candidate_page = ln
                break
        if candidate_page and not fields["페이지명"]:
            fields["페이지명"] = candidate_page

        # c) 검사항목: "숫자. (항목명) ..." 또는 괄호 안 항목명 패턴
        #    예) "6. (색에 무관한 콘텐츠 인식) 콘텐츠는 ..."
        item = ""
        for ln in lines:
            m = re.match(r"^\s*\d+\.\s*\(([^)]+)\)", ln)
            if m:
                item = m.group(1).strip()
                break
            m2 = re.search(r"\(([^)]+)\)\s*콘텐츠", ln)
            if m2 and "콘텐츠" in ln:
                item = m2.group(1).strip()
                break
        if item and not fields["검사항목"]:
            fields["검사항목"] = item

        # d) 문제점 및 개선 방안(텍스트):
        #    - 진단/지침 문장을 제외하고, URL/페이지명/메타성 문구를 뺀 나머지 서술을 모음
        diagnosis_headers = {"진단 결과 및 개선 방안", "동일 유형 페이지"}
        exclude_prefixes = (
            "http://", "https://",
        )
        exclude_exact = set([fields["페이지명"], fields["URL"]])
        collected: List[str] = []
        for ln in lines:
            if not ln:
                continue
            if ln in diagnosis_headers:
                continue
            if ln in exclude_exact:
                continue
            if ln.startswith(exclude_prefixes):
                continue
            # 이미 검사항목으로 인식한 줄은 제외
            if fields["검사항목"] and fields["검사항목"] in ln:
                continue
            collected.append(ln)

        # 짧은 잡음 제거: 헤더/라벨 느낌의 아주 짧은 줄 제거
        collected = [c for c in collected if len(c) >= 2]

        if collected and not fields["문제점 및 개선 방안(텍스트)"]:
            fields["문제점 및 개선 방안(텍스트)"] = "\n".join(collected)

        # 오류유형은 본문에 잘 안 나오는 경우가 많아 비워 둘 수 있음

    return fields


def _parse_fields_from_slide_tables(slide) -> Dict[str, str]:
    """
    Parse metadata fields by reading TABLE shapes directly.
    Assumptions:
    - First row of each table contains headers (labels).
    - Subsequent rows contain values.
    - Headers may include any alias in KOREAN_FIELD_ALIASES.
    - Multi-line text inside a cell is preserved as-is.
    If multiple rows/columns map to the same field, their content is concatenated with newlines.
    """
    fields: Dict[str, str] = {k: "" for k in FIELD_KEYS_ORDER}
    alias_to_field = {}
    for canon, aliases in KOREAN_FIELD_ALIASES.items():
        for a in aliases:
            alias_to_field[a.strip()] = canon

    def _norm(s: str) -> str:
        return s.replace("\r", "").strip()

    for shp, _, _ in _iter_shapes_recursive(slide):
        try:
            if shp.shape_type != MSO_SHAPE_TYPE.TABLE:
                continue
            tbl = shp.table
            if len(tbl.rows) < 1:
                continue
            headers = [_norm(c.text) for c in tbl.rows[0].cells]

            # Map columns to canonical field names
            col_to_field: Dict[int, str] = {}
            for idx, h in enumerate(headers):
                if not h:
                    continue
                if h in alias_to_field:
                    col_to_field[idx] = alias_to_field[h]
                    continue
                # loose match
                for alias, canon in alias_to_field.items():
                    if h.startswith(alias) or alias in h:
                        col_to_field[idx] = canon
                        break

            if not col_to_field:
                continue

            # Collect values
            for ridx in range(1, len(tbl.rows)):
                row = tbl.rows[ridx]
                for cidx, cell in enumerate(row.cells):
                    if cidx not in col_to_field:
                        continue
                    canon_key = col_to_field[cidx]
                    val = _norm(cell.text)
                    if not val:
                        continue
                    if fields[canon_key]:
                        fields[canon_key] += "\n" + val
                    else:
                        fields[canon_key] = val
        except Exception:
            continue

    return fields


def _build_eval_text(fields: Dict[str, str]) -> str:
    parts = []
    for k in FIELD_KEYS_ORDER:
        v = fields.get(k, "")
        if v:
            parts.append(f"{k}: {v}")
    return "\n".join(parts)


def generate_rationale(
    image_path: str,
    guideline_text: str,
    eval_text: str,
    model: str,
    api_key: str,
    dry_run: bool = False
) -> str:
    """Call an LLM to generate rationale. Skips if dry_run or no api_key."""
    if dry_run or not api_key:
        return ""

    client, mode = _lazy_openai_client(api_key)
    prompt = PROMPT_TEMPLATE.format(image_path=image_path, guide=guideline_text, eval=eval_text)

    try:
        if mode == "responses":
            # OpenAI Responses API
            resp = client.responses.create(
                model=model,
                input=[{"role": "user", "content": prompt}],
                temperature=0.2,
            )
            # Robust extraction
            text = getattr(resp, "output_text", None)
            if not text:
                # Fallback to content parts, if present
                content = getattr(resp, "content", None)
                if isinstance(content, list) and content:
                    # handle objects with .text attr or dict-like
                    first = content[0]
                    text = getattr(first, "text", None) or getattr(first, "content", None)
            return (text or "").strip()
        else:
            # Legacy Chat Completions
            resp = client.ChatCompletion.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.2,
            )
            return resp["choices"][0]["message"]["content"].strip()
    except Exception:
        # Don't fail the whole pipeline if LLM fails
        return ""


def process_pptx(
    pptx_path: str,
    out_root: str,
    model: str,
    api_key: str,
    guideline_text: str = "",
    dry_run: bool = False,
    debug: bool = False
) -> str:
    prs = Presentation(pptx_path)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    out_dir = os.path.join(out_root, base)
    os.makedirs(out_dir, exist_ok=True)

    jsonl_path = os.path.join(out_root, f"{base}.jsonl")
    with open(jsonl_path, "w", encoding="utf-8") as fw:
        for si, slide in enumerate(prs.slides, start=1):
            # 1) Extract images
            image_paths = _extract_images_from_slide(slide, out_dir, si)

            # 2) Extract texts & parse fields
            debug_fp = None
            if debug:
                dbgdir = os.path.join(out_root, "debug")
                os.makedirs(dbgdir, exist_ok=True)
                debug_fp = open(os.path.join(dbgdir, f"{base}_s{si:03d}.txt"), "w", encoding="utf-8")

            try:
                texts = _extract_text_from_slide(slide, debug_sink=debug_fp)
            finally:
                if debug_fp is not None:
                    debug_fp.close()

            # Prefer table-driven parsing; fallback to free-text parsing
            fields = _parse_fields_from_slide_tables(slide)
            if not any(fields.values()):
                fields = _parse_fields_from_texts(texts)

            # If nothing extracted, skip slide
            has_any = any(v for v in fields.values()) or len(image_paths) > 0
            if not has_any:
                continue

            eval_text = _build_eval_text(fields)

            # 3) Generate rationale per first image (if any). If no image, still attempt rationale.
            primary_img = image_paths[0] if image_paths else ""
            rationale = generate_rationale(
                image_path=primary_img,
                guideline_text=guideline_text,
                eval_text=eval_text,
                model=model,
                api_key=api_key,
                dry_run=dry_run,
            )

            # 4) Combine into JSON and write line
            record = {
                "pptx": os.path.basename(pptx_path),
                "slide_index": si,
                "image_paths": [os.path.relpath(p, out_root) for p in image_paths],
                "metadata": fields,
                "rationale": rationale,
            }
            fw.write(json.dumps(record, ensure_ascii=False) + "\n")

    return jsonl_path


def main():
    parser = argparse.ArgumentParser(description="Extract PPTX images/metadata and generate accessibility rationale JSONL.")
    parser.add_argument("pptx", help="Path to .pptx file")
    parser.add_argument("--out-dir", default="./train", help="Output root directory (default: ./train)")
    parser.add_argument("--model", default="gpt-4o", help="Reasoning model name (e.g., gpt-4o, gpt-4o-mini)")
    parser.add_argument("--openai_api_key", default=os.environ.get("OPENAI_API_KEY", ""), help="OpenAI API key or leave empty to skip")
    parser.add_argument("--guide-file", default="", help="Optional path to a guideline text/markdown file to feed into the prompt")
    parser.add_argument("--dry-run", action="store_true", help="Skip LLM calls and leave 'rationale' empty")
    parser.add_argument("--debug", action="store_true", help="Dump per-slide text extraction for debugging")
    args = parser.parse_args()

    os.makedirs(args.out_dir, exist_ok=True)

    guideline_text = ""
    if args.guide_file and os.path.exists(args.guide_file):
        with open(args.guide_file, "r", encoding="utf-8") as fr:
            guideline_text = fr.read()

    jsonl_path = process_pptx(
        pptx_path=args.pptx,
        out_root=args.out_dir,
        model=args.model,
        api_key=args.openai_api_key,
        guideline_text=guideline_text,
        dry_run=args.dry_run or not bool(args.openai_api_key),
        debug=args.debug,
    )

    print(f"Wrote JSONL to: {jsonl_path}")


if __name__ == "__main__":
    main()
