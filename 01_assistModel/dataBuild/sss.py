# recapture_metadata_playwright.py (visited deque 캐시 버전)
#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse, pathlib, sys, time, subprocess, json, re, os, io
from datetime import datetime
from urllib.parse import urlparse
from typing import Optional, List, Dict, Set
from collections import deque

from playwright.sync_api import sync_playwright, TimeoutError as PWTimeout
from pptx import Presentation

try:
    from PIL import Image
except Exception:
    Image = None

DEFAULT_HIDE_SELECTORS = [
    "nav", ".gnb", ".cookie", ".cookies",
    ".banner", ".popup", ".pop", ".floating", ".float",
    ".chat", ".chatbot", ".subscribe", ".sticky",
    ".toolbar", ".footer-quick", ".video-overlay"
]
HEADER_CANDIDATES = ['header', '.header', '#header', '[role="banner"]']

SLIDE_IDX_RE = re.compile(r"(?:^|/|\\)전체영역/(\d+)(?:_\d+)?\.png$", re.IGNORECASE)
URL_RE       = re.compile(r"https?://[^\s<>\"]+", re.IGNORECASE)

# ===== visited 캐시 =====
VISITED_MAX_DEFAULT = 256
VISITED_Q: deque[str] = deque(maxlen=VISITED_MAX_DEFAULT)
VISITED_MAP: Dict[str, bytes] = {}

def canon_url(url: str) -> str:
    """간단 정규화: 스킴/호스트 소문자, 트레일링 슬래시 제거."""
    p = urlparse(url.strip())
    host = (p.netloc or "").lower()
    scheme = (p.scheme or "").lower()
    path = p.path.rstrip("/") or "/"
    qs = ("?" + p.query) if p.query else ""
    frag = ""  # #fragment는 무시
    return f"{scheme}://{host}{path}{qs}{frag}"

def cache_get(url: str) -> Optional[bytes]:
    key = canon_url(url)
    data = VISITED_MAP.get(key)
    if data is not None:
        # LRU 갱신
        try:
            VISITED_Q.remove(key)
        except ValueError:
            pass
        VISITED_Q.append(key)
    return data

def cache_put(url: str, data: bytes, maxlen: Optional[int] = None):
    if maxlen is not None and VISITED_Q.maxlen != maxlen:
        # 런타임에 사이즈 바뀌면 재생성
        global VISITED_Q
        old = list(VISITED_Q)
        VISITED_Q = deque(old, maxlen=maxlen)
    key = canon_url(url)
    if key not in VISITED_MAP and len(VISITED_Q) == VISITED_Q.maxlen:
        # 제거
        old_key = VISITED_Q.popleft()
        VISITED_MAP.pop(old_key, None)
    else:
        try:
            VISITED_Q.remove(key)
        except ValueError:
            pass
    VISITED_Q.append(key)
    VISITED_MAP[key] = data

# ===== 유틸 =====
def guess_filename(url: str, ext: str) -> str:
    host = (urlparse(url).netloc or "page").replace(":", "_")
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    return f"{host}_{ts}.{ext}"

def quiet_install_chromium():
    try:
        subprocess.run(
            [sys.executable, "-m", "playwright", "install", "chromium", "--with-deps"],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )
    except Exception:
        pass

def add_stability_styles(page):
    css = """
    * { animation: none !important; transition: none !important; }
    html, body { scroll-behavior: auto !important; }
    """
    try: page.add_style_tag(content=css)
    except Exception: pass

# ===== 헤더/스크롤 유틸 =====
def mark_and_hide_header(page, candidates: List[str]):
    page.evaluate(
        r"""(args) => {
            const sels = args.sels;
            for (const sel of sels) {
                document.querySelectorAll(sel).forEach(el => {
                    el.setAttribute('data-fullshot-header','1');
                    el.style.setProperty('display','none','important');
                });
            }
        }""",
        {"sels": candidates},
    )

def show_header_static(page):
    page.evaluate(
        r"""() => {
            document.querySelectorAll('[data-fullshot-header]').forEach(el => {
                el.style.removeProperty('display');
                el.style.setProperty('position','static','important');
                el.style.setProperty('top','auto','important');
                el.style.setProperty('transform','none','important');
                el.querySelectorAll('*').forEach(child => {
                    const cs = getComputedStyle(child);
                    if (cs.position === 'sticky') {
                        child.style.setProperty('position','static','important');
                        child.style.setProperty('top','auto','important');
                    }
                });
            });
        }"""
    )

def header_exists(page) -> bool:
    try:
        return page.evaluate(r"""(cands) => {
            for (const sel of cands) {
                if (document.querySelector(sel)) return true;
            }
            return false;
        }""", HEADER_CANDIDATES)
    except Exception:
        return False

def hide_fixed_and_selectors(page, hide_selectors: List[str], max_h=220):
    try:
        page.evaluate(
            r"""(args) => {
                const sels = args.sels, maxH = args.maxH;
                for (const sel of sels) {
                    document.querySelectorAll(sel).forEach(el => {
                        if (el.closest('[data-fullshot-header]')) return;
                        el.style.setProperty('display','none','important');
                    });
                }
                const all = Array.from(document.querySelectorAll('body *')).slice(0, 5000);
                for (const el of all) {
                    if (el.closest('[data-fullshot-header]')) continue;
                    const cs = getComputedStyle(el);
                    if ((cs.position === 'fixed' || cs.position === 'sticky')) {
                        const r = el.getBoundingClientRect();
                        if (r.height > 0 && r.height <= maxH) {
                            el.style.setProperty('display','none','important');
                        }
                    }
                }
            }""",
            {"sels": hide_selectors, "maxH": int(max_h)}
        )
    except Exception:
        pass

def autoscroll_until_stable(page, step_px=900, idle_ms=800, stable_rounds=2, max_steps=500):
    last_h = 0; stable = 0; steps = 0
    while steps < max_steps:
        steps += 1
        try: page.evaluate("(s)=>window.scrollBy(0,s)", step_px)
        except Exception: break
        page.wait_for_timeout(50)
        try: page.wait_for_load_state("networkidle", timeout=idle_ms+700)
        except Exception: pass
        page.wait_for_timeout(idle_ms)
        try: h = page.evaluate("() => document.documentElement.scrollHeight")
        except Exception: break
        if h == last_h:
            stable += 1
            if stable >= stable_rounds: break
        else:
            stable = 0; last_h = h
    try:
        page.evaluate("() => window.scrollTo(0, document.documentElement.scrollHeight)")
        page.wait_for_timeout(idle_ms)
    except Exception: pass

# ===== 메모리 합성 =====
def compose_header_and_body_bytes(header_bytes: bytes, body_bytes: bytes, out_fmt: str, quality: Optional[int]) -> bytes:
    if Image is None:
        raise RuntimeError("Pillow가 필요합니다.  pip install pillow")
    header_img = Image.open(io.BytesIO(header_bytes)).convert("RGB")
    body_img   = Image.open(io.BytesIO(body_bytes)).convert("RGB")
    width = min(header_img.width, body_img.width)
    if header_img.width != width: header_img = header_img.crop((0,0,width,header_img.height))
    if body_img.width   != width: body_img   = body_img.crop((0,0,width,body_img.height))
    merged = Image.new("RGB", (width, header_img.height + body_img.height), (255,255,255))
    merged.paste(header_img, (0,0))
    merged.paste(body_img, (0, header_img.height))
    buff = io.BytesIO()
    if out_fmt.lower() in ("jpg","jpeg"):
        q = 90 if quality is None else max(1, min(100, quality))
        merged.save(buff, format="JPEG", quality=q)
    else:
        merged.save(buff, format="PNG")
    return buff.getvalue()

# ===== 단일 URL 캡처: bytes 반환 (캐시에 넣기 위함) =====
def capture_bytes(url: str, width: int, scale: float, timeout_s: int,
                  fmt: str, quality: Optional[int], channel: Optional[str],
                  include_header: bool, hide_selectors: List[str],
                  step_px: int, idle_ms: int) -> bytes:

    with sync_playwright() as p:
        try:
            browser = p.chromium.launch(headless=True)
        except Exception as e:
            if "Executable doesn't exist" in str(e):
                quiet_install_chromium()
                browser = p.chromium.launch(headless=True)
            else:
                browser = p.chromium.launch(headless=True, channel=(channel or "msedge"))

        context = browser.new_context(
            viewport={"width": width, "height": 800},
            device_scale_factor=scale,
            locale="ko-KR",
            timezone_id="Asia/Seoul"
        )
        page = context.new_page()
        add_stability_styles(page)

        page.goto(url, wait_until="domcontentloaded", timeout=timeout_s*1000)
        try: page.wait_for_load_state("networkidle", timeout=timeout_s*1000)
        except PWTimeout: pass

        if include_header:
            if header_exists(page):
                mark_and_hide_header(page, HEADER_CANDIDATES)
            autoscroll_until_stable(page, step_px=step_px, idle_ms=idle_ms)
            if hide_selectors:
                hide_fixed_and_selectors(page, hide_selectors, max_h=220)
            body_bytes = page.screenshot(full_page=True, type="png")

            show_header_static(page)
            page.evaluate("() => window.scrollTo(0, 0)")
            page.wait_for_timeout(idle_ms)
            try:
                header_el = page.locator("[data-fullshot-header]").first
                if header_el.count() == 0:
                    for sel in HEADER_CANDIDATES:
                        if page.locator(sel).count() > 0:
                            header_el = page.locator(sel).first
                            break
                header_bytes = header_el.screenshot(type="png")
                final_bytes = compose_header_and_body_bytes(header_bytes, body_bytes, out_fmt=fmt, quality=quality)
            except Exception:
                final_bytes = page.screenshot(full_page=True, type=("jpeg" if fmt=="jpeg" else "png"))
        else:
            if header_exists(page):
                mark_and_hide_header(page, HEADER_CANDIDATES)
            autoscroll_until_stable(page, step_px=step_px, idle_ms=idle_ms)
            if hide_selectors:
                hide_fixed_and_selectors(page, hide_selectors, max_h=220)
            final_bytes = page.screenshot(full_page=True, type=("jpeg" if fmt=="jpeg" else "png"))

        context.close(); browser.close()
        return final_bytes

# ===== PPTX → 슬라이드별 URL =====
def extract_urls_by_slide(pptx_path: str) -> Dict[int, str]:
    if not os.path.exists(pptx_path):
        print(f"[ERR] PPTX not found: {pptx_path}")
        return {}
    prs = Presentation(pptx_path)
    slide_to_url: Dict[int, str] = {}
    for i, slide in enumerate(prs.slides, start=1):
        urls: List[str] = []
        for shp in slide.shapes:
            try:
                link = getattr(getattr(shp, "click_action", None), "hyperlink", None)
                if link and link.address: urls.append(link.address)
            except Exception: pass
            try:
                link2 = getattr(shp, "hyperlink", None)
                if link2 and link2.address: urls.append(link2.address)
            except Exception: pass
            if getattr(shp, "has_text_frame", False):
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        try:
                            if r.hyperlink and r.hyperlink.address:
                                urls.append(r.hyperlink.address)
                        except Exception: pass
        texts=[]
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                for p in shp.text_frame.paragraphs:
                    texts.append("".join(r.text or "" for r in p.runs))
            if getattr(shp, "has_table", False):
                for r in shp.table.rows:
                    for c in r.cells:
                        if c.text: texts.append(c.text)
        blob = "\n".join(texts)
        for m in URL_RE.finditer(blob): urls.append(m.group(0))
        seen: Set[str] = set()
        for u in urls:
            u2 = u.strip("<> ,.;)]}'\"")
            if u2.lower().startswith(("http://","https://")) and u2 not in seen:
                slide_to_url.setdefault(i, u2); seen.add(u2)
    return slide_to_url

# ===== metadata.jsonl 기반 일괄 재캡처 (+ visited 캐시 적용) =====
def recapture_from_metadata(pptx_path: str, metadata_path: str,
                            width: int, scale: float, timeout_s: int,
                            include_header: bool, hide_selectors: List[str],
                            step_px: int, idle_ms: int, fmt: str, quality: Optional[int],
                            channel: Optional[str], cache_size: int):
    print("[DBG] recapture start")
    print("  PPTX :", os.path.abspath(pptx_path))
    print("  META :", os.path.abspath(metadata_path))
    # 캐시 크기 설정
    cache_put("about:blank", b"", maxlen=cache_size)  # maxlen 반영용 더미
    VISITED_MAP.pop(canon_url("about:blank"), None)
    try:
        VISITED_Q.remove(canon_url("about:blank"))
    except ValueError:
        pass

    slide_urls = extract_urls_by_slide(pptx_path)
    if not os.path.exists(metadata_path):
        print(f"[ERR] metadata.jsonl not found: {metadata_path}")
        return

    processed = 0
    with open(metadata_path, "r", encoding="utf-8") as f:
        for ln, line in enumerate(f, start=1):
            s = line.strip()
            if not s: continue
            try:
                rec = json.loads(s)
            except Exception as e:
                print(f"[WARN] line {ln} json error: {e}")
                continue

            targets = [p for p in (rec.get("file_names", []) or []) if "전체영역" in p]
            for out_path in targets:
                m = SLIDE_IDX_RE.search(out_path)
                if not m:
                    print(f"[SKIP] l{ln}: slide idx not found in {out_path}")
                    continue
                idx = int(m.group(1))
                url = slide_urls.get(idx)
                if not url:
                    print(f"[SKIP] slide {idx}: url not found in pptx")
                    continue

                print(f"[GO] slide {idx} -> {url} -> {out_path}")
                os.makedirs(os.path.dirname(out_path), exist_ok=True)

                # === 캐시 HIT 체크 ===
                cached = cache_get(url)
                if cached is not None:
                    pathlib.Path(out_path).write_bytes(cached)
                    print(f"[HIT] cache used for {url} ({len(cached)} bytes)")
                    processed += 1
                    continue

                # === 캡처 수행 ===
                try:
                    img_bytes = capture_bytes(
                        url=url, width=width, scale=scale, timeout_s=timeout_s,
                        fmt=fmt, quality=quality, channel=channel,
                        include_header=include_header, hide_selectors=hide_selectors,
                        step_px=step_px, idle_ms=idle_ms
                    )
                    pathlib.Path(out_path).write_bytes(img_bytes)
                    cache_put(url, img_bytes, maxlen=cache_size)
                    processed += 1
                    print(f"[OK] saved: {out_path} ({len(img_bytes)} bytes)")
                except PWTimeout:
                    print(f"[TIMEOUT] slide {idx} @ {url}")
                except Exception as e:
                    print(f"[ERR] slide {idx}: {e}")

    print(f"[DONE] processed: {processed}")

def main():
    ap = argparse.ArgumentParser(description="metadata.jsonl + PPTX 기반 전체영역 재캡처 (Playwright, visited deque 캐시)")
    ap.add_argument("--pptx", default="sample.pptx")
    ap.add_argument("--metadata", default="metadata.jsonl")
    ap.add_argument("--width", type=int, default=1366)
    ap.add_argument("--scale", type=float, default=2.0)
    ap.add_argument("--timeout", type=int, default=60)
    ap.add_argument("--fmt", choices=["png","jpeg"], default="png")
    ap.add_argument("--quality", type=int)
    ap.add_argument("--channel", choices=["msedge","chrome"])
    ap.add_argument("--include-header", action="store_true")
    ap.add_argument("--hide-selectors", default=",".join(DEFAULT_HIDE_SELECTORS))
    ap.add_argument("--scroll-step", type=int, default=1200)
    ap.add_argument("--idle-ms", type=int, default=1200)
    ap.add_argument("--cache-size", type=int, default=VISITED_MAX_DEFAULT,
                    help="visited deque 캐시의 최대 URL 수 (기본 256)")
    args = ap.parse_args()

    hide_sels = [s.strip() for s in (args.hide_selectors or "").split(",") if s.strip()]

    recapture_from_metadata(
        pptx_path=args.pptx, metadata_path=args.metadata,
        width=args.width, scale=args.scale, timeout_s=args.timeout,
        include_header=args.include_header, hide_selectors=hide_sels,
        step_px=args.scroll_step, idle_ms=args.idle_ms,
        fmt=args.fmt, quality=args.quality, channel=args.channel,
        cache_size=args.cache_size
    )

if __name__ == "__main__":
    main()
